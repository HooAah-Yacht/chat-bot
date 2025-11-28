"""
HooAah Yacht AI Chatbot - 통합 버전
모든 챗봇 기능을 하나로 통합한 통합 챗봇

기능:
- 자연어 대화 (Gemini AI)
- PDF 업로드 및 분석 (자동 스펙 및 부품 추출)
- 기존 20종 요트 정보 조회
- API 서버 모드 (Flutter 앱 연동)
- 자동 명령어 인식 및 처리
- JSON 파일 자동 저장

사용법:
    # 대화형 모드
    python chatbot_unified.py
    
    # API 서버 모드
    python chatbot_unified.py --mode api
    
    # API 키 지정
    python chatbot_unified.py --api-key YOUR_API_KEY
"""

import os
import sys

# Windows 콘솔 인코딩 설정
if sys.platform == 'win32':
    try:
        sys.stdout.reconfigure(encoding='utf-8')
        sys.stderr.reconfigure(encoding='utf-8')
    except:
        pass
import json
import sys
import argparse
from datetime import datetime
from typing import List, Dict, Optional
from pathlib import Path

# Environment variables (.env 파일 로드)
try:
    from dotenv import load_dotenv
    load_dotenv()  # .env 파일에서 환경 변수 로드
    HAS_DOTENV = True
except ImportError:
    HAS_DOTENV = False
    try:
        print("[WARNING] python-dotenv 패키지가 설치되지 않았습니다. pip install python-dotenv")
    except:
        print("[WARNING] python-dotenv package not installed. pip install python-dotenv")

# Gemini AI 관련
try:
    import google.generativeai as genai
    HAS_GEMINI = True
except ImportError:
    HAS_GEMINI = False
    try:
        print("[WARNING] google-generativeai 패키지가 설치되지 않았습니다. pip install google-generativeai")
    except:
        print("[WARNING] google-generativeai package not installed. pip install google-generativeai")

# Flask API 관련
try:
    from flask import Flask, request, jsonify
    from flask_cors import CORS
    from werkzeug.utils import secure_filename
    HAS_FLASK = True
except ImportError:
    HAS_FLASK = False
    secure_filename = None
    try:
        print("[WARNING] flask 패키지가 설치되지 않았습니다.")
        print("[INFO] 자동 설치를 시도합니다...")
    except:
        print("[WARNING] flask package not installed.")
        print("[INFO] Attempting auto-install...")
    import subprocess
    import sys
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "flask", "flask-cors", "--quiet"])
        from flask import Flask, request, jsonify
        from flask_cors import CORS
        from werkzeug.utils import secure_filename
        HAS_FLASK = True
        try:
            print("[SUCCESS] Flask 패키지 설치 완료!")
        except:
            print("[SUCCESS] Flask package installed!")
    except Exception as e:
        try:
            print(f"[ERROR] Flask 자동 설치 실패: {e}")
            print("[INFO] 수동 설치: pip install flask flask-cors")
        except:
            print(f"[ERROR] Flask auto-install failed: {e}")
            print("[INFO] Manual install: pip install flask flask-cors")

# PDF 분석 관련
try:
    import PyPDF2
    HAS_PYPDF2 = True
except ImportError:
    HAS_PYPDF2 = False

try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False

# OCR 관련 (선택사항)
try:
    import pytesseract
    from pdf2image import convert_from_path
    HAS_OCR = True
except ImportError:
    HAS_OCR = False

# Word 문서 처리
try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

# HWP 파일 처리
try:
    import olefile
    HAS_OLEFILE = True
except ImportError:
    HAS_OLEFILE = False

# Excel 파일 처리
try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

# PowerPoint 파일 처리
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class UnifiedYachtChatbot:
    """
    통합 요트 챗봇 클래스
    모든 챗봇 기능을 통합
    """
    
    def __init__(self, api_key: str = None, mode: str = "interactive"):
        """
        통합 챗봇 초기화
        
        Args:
            api_key: Gemini API 키 (기본값: 제공된 API 키)
            mode: 실행 모드 ("interactive", "api", "cli")
        """
        self.mode = mode
        # API 키 우선순위: 인자 > 환경변수 (.env 파일) > 없으면 오류
        self.api_key = api_key or os.getenv('GEMINI_API_KEY')
        if not self.api_key:
            print("⚠️ 경고: GEMINI_API_KEY가 설정되지 않았습니다.")
            print("💡 .env 파일에 GEMINI_API_KEY를 설정하거나 --api-key 옵션을 사용하세요.")
            print("📝 .env.example 파일을 참고하여 .env 파일을 생성하세요.")
        
        # Gemini AI 초기화
        if HAS_GEMINI and self.api_key:
            genai.configure(api_key=self.api_key)
            try:
                self.model = genai.GenerativeModel('gemini-2.5-flash')
                print("✅ Gemini 2.5 Flash 모델 사용")
            except Exception as e:
                print(f"⚠️ Gemini 2.5 Flash 사용 실패, gemini-pro로 전환: {e}")
                self.model = genai.GenerativeModel('gemini-pro')
                print("✅ gemini-pro 모델 사용 (fallback)")
            self.has_gemini = True
        else:
            self.has_gemini = False
            print("⚠️ Gemini AI를 사용할 수 없습니다. 기본 모드로 실행됩니다.")
        
        # 대화 히스토리
        self.chat_history: List[Dict[str, str]] = []
        
        # 요트 데이터 로드
        self.yacht_data = self._load_yacht_data()
        self.parts_data = self._load_parts_data()
        
        # 시스템 프롬프트
        if self.has_gemini:
            self.system_prompt = self._create_system_prompt()
        
        # 등록 데이터
        self.current_yacht_registration = None
        
        print("✅ HooAah Yacht 통합 챗봇이 준비되었습니다!")
        if mode == "interactive":
            print("💬 자연스럽게 요트에 대해 질문해보세요.")
            print("📄 PDF 파일 경로를 입력하면 자동으로 분석합니다.\n")
    
    def _load_yacht_data(self) -> Dict:
        """요트 스펙 데이터 로드"""
        try:
            with open('data/yacht_specifications.json', 'r', encoding='utf-8') as f:
                return json.load(f)
        except FileNotFoundError:
            print("⚠️ yacht_specifications.json 파일을 찾을 수 없습니다.")
            return {"yachts": []}
    
    def _load_parts_data(self) -> Dict:
        """부품 데이터 로드 (interval 정보가 있는 yacht_parts_app_data.json 우선)"""
        try:
            # yacht_parts_app_data.json을 우선 로드 (interval 정보 포함)
            with open('data/yacht_parts_app_data.json', 'r', encoding='utf-8') as f:
                return json.load(f)
        except FileNotFoundError:
            # fallback: yacht_parts_database.json
            try:
                with open('data/yacht_parts_database.json', 'r', encoding='utf-8') as f:
                    return json.load(f)
            except FileNotFoundError:
                print("⚠️ 부품 데이터 파일을 찾을 수 없습니다.")
                return {"yachts": []}
    
    def _create_system_prompt(self) -> str:
        """시스템 프롬프트 생성"""
        yacht_list = [yacht.get('name', '') for yacht in self.yacht_data.get('yachts', [])]
        
        prompt = f"""당신은 HooAah Yacht의 전문 AI 어시스턴트입니다.

**역할:**
- 요트 소유자와 관리자를 돕는 친절하고 전문적인 어시스턴트
- 요트 스펙, 부품, 정비, 관리에 대한 모든 질문에 답변
- 자연스럽고 대화적인 톤으로 소통
- 사용자의 의도를 파악하여 적절한 응답 제공

**지원하는 요트 20종:**
{', '.join(yacht_list)}

**데이터베이스 구조 (ERD 기반):**
- User: 사용자 정보
- Yacht: 요트 정보 (name)
- Yacht_User: 사용자-요트 연결 (다대다 관계)
- Part: 부품 정보 (name, manufacturer, model, interval)
- Repair: 정비 내역 (repairDate만 저장, content 필드 없음)
- Calendar: 캘린더 이벤트 (content 필드 있음, part_id와 연결)
- Alert: 알림 (part_id와 일대일 관계)

**의도 파악 및 응답 가이드라인:**

1. **요트 정보 조회 의도:**
   - 키워드: "정보", "스펙", "사양", "크기", "치수", "길이", "폭", "높이", 요트 이름 등
   - 응답: 해당 요트의 상세 정보 제공 (치수, 엔진, 돛 면적 등)

2. **요트 등록/PDF 업로드 의도:**
   - 키워드: "등록", "업로드", "pdf", "문서", "매뉴얼", "새 요트", "요트 등록", "추가", "입력" 등
   - 응답: PDF 파일 업로드 안내 메시지 반환
   - 예시 문장:
     * "새 요트 등록하고 싶어요"
     * "pdf 매뉴얼 업로드해줘"
     * "요트 정보 추가할게"
     * "문서 등록하고 싶어"
     * "매뉴얼 파일 올릴 수 있어?"

3. **도움말 요청 의도:**
   - 키워드: "도움말", "도움", "help", "사용법", "어떻게", "방법", "사용", "사용법", "가이드" 등
   - 응답: 도움말 메시지 반환

4. **요트 목록 조회 의도:**
   - 키워드: "목록", "리스트", "전체", "모든 요트", "어떤 요트", "요트 종류" 등
   - 응답: 요트 목록 반환

5. **요트 비교 의도:**
   - 키워드: "비교", "차이", "어떤 게", "vs", "대" 등
   - 응답: 여러 요트를 비교하여 차이점 설명

6. **요트 추천 의도:**
   - 키워드: "추천", "어떤 게 좋아", "선택", "고르" 등
   - 응답: 사용 목적에 맞는 요트 추천

7. **정비/관리 질문:**
   - 키워드: "정비", "관리", "주기", "점검", "교체" 등
   - 응답: 정비 주기, 정비 이력, 관리 방법 안내

**답변 형식:**
- 짧고 명확하게 (모바일 화면에 적합)
- 필요시 이모지 사용 (⛵, 🔧, 📏, ⚓ 등)
- 숫자는 단위와 함께 명시
- 추가 질문 유도

**중요:**
- 사용자의 의도를 정확히 파악하여 적절한 응답 제공
- 자연어로 된 모든 요청을 이해하고 처리
- 특정 명령어가 아닌 관련 키워드만 있어도 의도 파악
- 모르는 내용은 솔직히 모른다고 답변

**데이터 활용:**
아래 JSON 데이터를 참고하여 정확한 정보 제공:

요트 스펙 데이터:
{json.dumps(self.yacht_data, ensure_ascii=False, indent=2)[:5000]}...

부품 데이터 (샘플):
{json.dumps(self.parts_data, ensure_ascii=False, indent=2)[:3000]}...
"""
        return prompt
    
    def _get_file_extension(self, file_path: str) -> str:
        """파일 확장자 추출"""
        return os.path.splitext(file_path)[1].lower()
    
    def _is_supported_file(self, file_path: str) -> bool:
        """지원되는 파일 형식인지 확인"""
        ext = self._get_file_extension(file_path)
        supported_extensions = ['.pdf', '.docx', '.doc', '.hwp', '.txt', '.xlsx', '.xls', '.pptx', '.ppt']
        return ext in supported_extensions
    
    def _extract_file_path_from_message(self, message: str) -> Optional[str]:
        """
        사용자 메시지에서 파일 경로 추출 (PDF, Word, HWP 등)
        모바일 앱에서 전달된 파일 경로도 지원 (iOS, Android)
        """
        import re
        
        # 지원되는 파일 확장자
        supported_exts = r'(?:pdf|docx?|hwp|txt|xlsx?)'
        
        # 1. 따옴표로 감싸진 경로 찾기 (공백 포함 경로 지원)
        quoted_patterns = [
            rf'["\']([^"\']+\.{supported_exts})["\']',
            rf'["\']([^"\']+\.{supported_exts})',
            rf'([^"\']+\.{supported_exts})["\']',
        ]
        
        for pattern in quoted_patterns:
            match = re.search(pattern, message, re.IGNORECASE)
            if match:
                path = match.group(1).strip()
                if os.path.exists(path):
                    return os.path.abspath(path)
        
        # 2. Windows 절대 경로 패턴
        windows_abs_pattern = rf'([A-Za-z]:[\\/](?:[^"\']+[\\/])*[^"\']+\.{supported_exts})'
        match = re.search(windows_abs_pattern, message, re.IGNORECASE)
        if match:
            path = match.group(1).strip()
            if os.path.exists(path):
                return os.path.abspath(path)
        
        # 3. Unix/Linux/Mac 절대 경로 패턴
        unix_abs_pattern = rf'(/[^"\']+\.{supported_exts})'
        match = re.search(unix_abs_pattern, message, re.IGNORECASE)
        if match:
            path = match.group(1).strip()
            if os.path.exists(path):
                return os.path.abspath(path)
        
        # 4. 모바일 앱 경로 패턴
        mobile_patterns = [
            rf'(/storage/[^"\']+\.{supported_exts})',  # Android
            rf'(/var/mobile/[^"\']+\.{supported_exts})',  # iOS
            rf'(/data/[^"\']+\.{supported_exts})',  # Android data
            rf'(file://[^"\']+\.{supported_exts})',  # file:// URI
        ]
        
        for pattern in mobile_patterns:
            match = re.search(pattern, message, re.IGNORECASE)
            if match:
                path = match.group(1).strip()
                if path.startswith('file://'):
                    path = path.replace('file://', '')
                if os.path.exists(path):
                    return os.path.abspath(path)
        
        # 5. 메시지 전체가 파일 경로인지 확인
        message_clean = message.strip().strip('"').strip("'")
        
        if os.path.isabs(message_clean) and self._is_supported_file(message_clean):
            if os.path.exists(message_clean):
                return os.path.abspath(message_clean)
        
        if self._is_supported_file(message_clean):
            if os.path.exists(message_clean):
                return os.path.abspath(message_clean)
            abs_path = os.path.abspath(message_clean)
            if os.path.exists(abs_path):
                return abs_path
        
        return None
    
    def _extract_pdf_path_from_message(self, message: str) -> Optional[str]:
        """메시지에서 PDF 파일 경로 추출 (하위 호환성)"""
        return self._extract_file_path_from_message(message)
    
    def _is_pdf_upload_request(self, message: str) -> bool:
        """파일 업로드 요청인지 확인 (PDF, Word, HWP 등)"""
        file_keywords = ['pdf', 'docx', 'doc', 'hwp', 'word', '문서', '매뉴얼', '업로드', '등록', '파일']
        return any(keyword in message.lower() for keyword in file_keywords) or \
               self._extract_file_path_from_message(message) is not None
    
    def _is_registration_request(self, message: str) -> bool:
        """요트 등록 요청인지 확인"""
        message_lower = message.lower()
        registration_keywords = [
            '요트 등록', '등록하고 싶어', '등록하고 싶어요', '등록하고 싶습니다',
            '새 요트', '요트 추가', '추가하고 싶어', '추가하고 싶어요',
            '부품 추가', '부품 등록', '부품 넣어', '부품 넣어줘'
        ]
        return any(keyword in message_lower for keyword in registration_keywords)
    
    def _handle_registration_request(self, user_message: str) -> str:
        """요트 등록/부품 추가 요청 처리"""
        message_lower = user_message.lower()
        
        # 부품 추가 요청인지 확인
        parts_keywords = ['부품 추가', '부품 등록', '부품 넣어', '부품 넣어줘']
        if any(keyword in message_lower for keyword in parts_keywords):
            # 기존 요트에 부품 추가
            yacht_name = self._extract_yacht_name_from_message(user_message)
            if yacht_name:
                return f"""📦 **{yacht_name} 부품 추가**

부품을 추가하려면 다음 방법 중 하나를 선택하세요:

1. **PDF 매뉴얼 업로드** (권장)
   - PDF 파일 경로를 입력하세요
   - 예: `data/yachtpdf/manual.pdf`

2. **수동 입력** (준비 중)
   - 곧 지원 예정입니다

PDF 파일 경로를 입력해주세요! 📎"""
            else:
                return """📦 **부품 추가**

어떤 요트에 부품을 추가하시겠어요?

1. 요트 이름을 알려주세요
   예: "Farr 40 부품 추가"

2. PDF 매뉴얼 파일 경로를 입력하세요
   예: `data/yachtpdf/manual.pdf`

요트 이름 또는 PDF 파일 경로를 입력해주세요! 📎"""
        
        # 일반 요트 등록 요청
        return self._suggest_pdf_upload()
    
    def _is_yacht_info_request(self, message: str) -> bool:
        """요트 정보 요청인지 확인"""
        info_keywords = ['정보', '스펙', '사양', '크기', '치수', '길이', '폭', '높이', '추천']
        yacht_names = [y.get('name', '').lower() for y in self.yacht_data.get('yachts', [])]
        
        message_lower = message.lower()
        return any(keyword in message_lower for keyword in info_keywords) or \
               any(name in message_lower for name in yacht_names)
    
    def chat(self, user_message: str, pdf_file_path: str = None) -> str:
        """
        사용자 메시지에 대한 응답 생성
        
        Args:
            user_message: 사용자 입력 메시지
            pdf_file_path: PDF 파일 경로 (모바일 앱에서 직접 전달되는 경우)
            
        Returns:
            AI 응답 메시지
        """
        try:
            # 1. 직접 전달된 파일 경로 확인 (모바일 앱에서 파일 업로드)
            if pdf_file_path and os.path.exists(pdf_file_path):
                return self._handle_file_upload(pdf_file_path)
            
            # 2. 메시지에서 PDF 파일 경로 추출
            pdf_path = self._extract_pdf_path_from_message(user_message)
            if pdf_path and os.path.exists(pdf_path):
                # 즉시 피드백
                print(f"\n📥 파일을 인식했습니다: {os.path.basename(pdf_path)}")
                print("⏳ 분석을 시작합니다. 잠시만 기다려주세요...\n")
                sys.stdout.flush()
                return self._handle_file_upload(pdf_path)
            
            # 3. 명령어 처리 (빠른 응답)
            message_lower = user_message.lower().strip()
            if message_lower in ['/list', '/목록']:
                return self._list_yachts()
            
            if message_lower in ['/info', '/정보']:
                return self._get_data_info()
            
            if message_lower in ['/help', '/도움말']:
                return self._get_help()
            
            # 4. 대화 히스토리에 추가
            self.chat_history.append({
                "role": "user",
                "content": user_message,
                "timestamp": datetime.now().isoformat()
            })
            
            # 5. 간단한 질문 먼저 처리 (크기, 부품 개수 등)
            simple_response = self._handle_simple_questions(user_message)
            if simple_response:
                response = simple_response
            # 5-1. 요트 등록/부품 추가 요청 처리
            elif self._is_registration_request(user_message):
                response = self._handle_registration_request(user_message)
            # 6. Gemini AI로 의도 파악 및 응답 생성
            elif self.has_gemini:
                # AI 응답 생성 시작 표시 (즉시)
                print("🤖 AI가 생각 중입니다...", end="", flush=True)
                # 분석 요청인 경우 추가 메시지
                if any(keyword in user_message.lower() for keyword in ['분석', '분석해줘', '분석해주세요', '상세 분석']):
                    print(" (상세 분석 중) ⏳", flush=True)
                else:
                    print(" ⏳", flush=True)
                # AI가 의도를 파악하여 적절한 응답 생성
                response = self._generate_intelligent_response(user_message)
                # 완료 표시 (줄바꿈)
                print("\r" + " " * 80 + "\r", end="", flush=True)  # 이전 메시지 지우기
            else:
                # 기본 모드: 키워드 기반 응답
                response = self._generate_keyword_based_response(user_message)
            
            # 7. 대화 히스토리에 추가
            self.chat_history.append({
                "role": "assistant",
                "content": response,
                "timestamp": datetime.now().isoformat()
            })
            
            return response
            
        except Exception as e:
            error_msg = f"죄송합니다. 응답 생성 중 오류가 발생했습니다: {str(e)}"
            print(f"❌ Error: {e}")
            return error_msg
    
    def _generate_intelligent_response(self, user_message: str) -> str:
        """Gemini AI를 사용한 지능형 응답 생성 (의도 파악)"""
        try:
            # 의도 파악을 위한 프롬프트
            intent_prompt = f"""사용자 메시지: "{user_message}"

위 메시지를 분석하여 사용자의 의도를 파악하고 적절한 응답을 생성해주세요.

**의도 분류:**
1. **요트 정보 조회**: 요트 이름, 스펙, 치수 등에 대한 질문
2. **요트 등록/PDF 업로드**: 새 요트를 등록하거나 PDF 매뉴얼을 업로드하려는 의도
3. **요트 분석 요청**: 기존 요트 데이터를 분석하거나 상세 분석을 요청하는 의도 (키워드: "분석", "분석해줘", "분석해주세요", "상세 분석", "데이터 분석" 등)
4. **도움말 요청**: 사용법, 가이드, 도움말을 요청하는 의도
5. **요트 목록 조회**: 전체 요트 목록을 보려는 의도
6. **요트 비교/추천**: 여러 요트를 비교하거나 추천을 요청하는 의도
7. **정비/관리 질문**: 정비 주기, 관리 방법 등에 대한 질문
8. **일반 대화**: 기타 요트 관련 질문

**응답 규칙:**
- 요트 등록/PDF 업로드 의도가 감지되면: PDF 파일 업로드 안내 메시지를 반환
- 요트 분석 요청 의도가 감지되면: 해당 요트의 상세 분석 정보를 제공 (스펙, 부품, 정비 주기, 특징 등 종합 분석)
- 도움말 요청 의도가 감지되면: 도움말 내용을 반환
- 요트 목록 조회 의도가 감지되면: 요트 목록을 반환
- 요트 정보 조회 의도가 감지되면: 해당 요트의 상세 정보를 제공
- 그 외: 자연스럽게 답변

**지원하는 요트 20종:**
{', '.join([yacht.get('name', '') for yacht in self.yacht_data.get('yachts', [])])}

위 규칙에 따라 사용자에게 적절한 응답을 생성해주세요."""
            
            # 의도 파악 및 응답 생성
            response = self.model.generate_content(intent_prompt)
            ai_response = response.text.strip()
            
            # 특수 응답 처리 (PDF 업로드, 분석, 도움말 등)
            ai_response_lower = ai_response.lower()
            user_message_lower = user_message.lower()
            
            # PDF 업로드 의도가 명확한 경우
            if any(keyword in ai_response_lower for keyword in ['pdf', '업로드', '등록', '파일 경로']):
                # PDF 업로드 안내 메시지로 대체
                return self._suggest_pdf_upload()
            
            # 요트 분석 의도가 명확한 경우
            if any(keyword in user_message_lower for keyword in ['분석', '분석해줘', '분석해주세요', '상세 분석', '데이터 분석', '요트 분석']):
                # 요트 이름이 포함되어 있으면 해당 요트 분석, 없으면 전체 분석 안내
                yacht_name = self._extract_yacht_name_from_message(user_message)
                if yacht_name:
                    # 특정 섹션 요청 감지 (예: "사용 목적", "적합성 평가", "정비 권장사항")
                    section_filter = self._extract_section_keyword(user_message)
                    return self._analyze_yacht_data(yacht_name, section_filter)
                else:
                    return "어떤 요트를 분석하시겠어요? 요트 이름을 알려주시면 상세 분석을 제공해드리겠습니다.\n예: 'Farr 40 분석해줘'"
            
            # 도움말 의도가 명확한 경우
            if any(keyword in ai_response_lower for keyword in ['도움말', '사용법', '가이드']):
                # 도움말 메시지로 대체
                return self._get_help()
            
            # 요트 목록 의도가 명확한 경우
            if any(keyword in ai_response_lower for keyword in ['목록', '리스트', '전체 요트']):
                # 요트 목록으로 대체
                return self._list_yachts()
            
            # 일반 응답 반환
            return ai_response
            
        except Exception as e:
            # 오류 발생 시 키워드 기반 응답으로 fallback
            return self._generate_keyword_based_response(user_message)
    
    def _generate_keyword_based_response(self, user_message: str) -> str:
        """키워드 기반 응답 생성 (Gemini AI 없을 때)"""
        message_lower = user_message.lower()
        
        # 1. 도움말 관련 키워드
        help_keywords = ['도움말', '도움', 'help', '사용법', '어떻게', '방법', '사용', '가이드', '안내']
        if any(keyword in message_lower for keyword in help_keywords):
            return self._get_help()
        
        # 2. PDF 업로드/등록 관련 키워드
        pdf_keywords = ['pdf', '문서', '매뉴얼', '업로드', '등록', '파일', '새 요트', '요트 등록', '추가', '입력', '올리', '넣']
        if any(keyword in message_lower for keyword in pdf_keywords):
            return self._suggest_pdf_upload_without_ai()
        
        # 2-1. 요트 분석 관련 키워드
        analysis_keywords = ['분석', '분석해줘', '분석해주세요', '상세 분석', '데이터 분석', '요트 분석']
        if any(keyword in message_lower for keyword in analysis_keywords):
            yacht_name = self._extract_yacht_name_from_message(user_message)
            if yacht_name:
                # 특정 섹션 요청 감지
                section_filter = self._extract_section_keyword(user_message)
                return self._analyze_yacht_data(yacht_name, section_filter)
            else:
                return "어떤 요트를 분석하시겠어요? 요트 이름을 알려주시면 상세 분석을 제공해드리겠습니다.\n예: 'Farr 40 분석해줘'"
        
        # 3. 요트 목록 관련 키워드
        list_keywords = ['목록', '리스트', '전체', '모든 요트', '어떤 요트', '요트 종류', '요트 목록']
        if any(keyword in message_lower for keyword in list_keywords):
            return self._list_yachts()
        
        # 4. 요트 정보 조회 (기존 로직)
        return self._generate_basic_response(user_message)
    
    def _generate_ai_response(self, user_message: str) -> str:
        """Gemini AI를 사용한 응답 생성"""
        try:
            context = self._build_context()
            response = self.model.generate_content(context)
            return response.text
        except Exception as e:
            return f"AI 응답 생성 중 오류가 발생했습니다: {str(e)}"
    
    def _generate_basic_response(self, user_message: str) -> str:
        """기본 모드 응답 생성 (Gemini AI 없이)"""
        # 기존 chatbot.py의 로직 사용
        message_lower = user_message.lower()
        
        # 요트 이름 찾기
        for yacht in self.yacht_data.get('yachts', []):
            yacht_name = yacht.get('name', '').lower()
            if yacht_name in message_lower:
                return self._format_yacht_info(yacht)
        
        return "죄송합니다. 요트 정보를 찾을 수 없습니다. '/list' 명령어로 요트 목록을 확인하세요."
    
    def _handle_simple_questions(self, user_message: str) -> Optional[str]:
        """간단한 질문 처리 (크기, 부품 개수 등)"""
        message_lower = user_message.lower()
        
        # 요트 이름 추출
        yacht_name = self._extract_yacht_name_from_message(user_message)
        if not yacht_name:
            return None
        
        # 요트 찾기
        yacht = None
        for y in self.yacht_data.get('yachts', []):
            if y.get('name', '').lower() == yacht_name.lower():
                yacht = y
                break
        
        if not yacht:
            return None
        
        # 1. 특정 치수 요소 질문 (개별 처리)
        # 폭 (Beam)
        beam_keywords = ['폭', 'beam', '너비', '가로']
        if any(keyword in message_lower for keyword in beam_keywords):
            return self._format_specific_dimension(yacht, 'beam', '폭 (Beam)')
        
        # 길이/전장 (LOA)
        loa_keywords = ['전장', 'loa', '길이', '전체 길이', '총 길이']
        if any(keyword in message_lower for keyword in loa_keywords):
            return self._format_specific_dimension(yacht, 'loa', '전장 (LOA)')
        
        # 흘수 (Draft)
        draft_keywords = ['흘수', 'draft', '드래프트']
        if any(keyword in message_lower for keyword in draft_keywords):
            return self._format_specific_dimension(yacht, 'draft', '흘수 (Draft)')
        
        # 배수량 (Displacement)
        displacement_keywords = ['배수량', 'displacement', '무게', '중량']
        if any(keyword in message_lower for keyword in displacement_keywords):
            return self._format_specific_dimension(yacht, 'displacement', '배수량 (Displacement)')
        
        # 마스트 높이
        mast_keywords = ['마스트', 'mast', '마스트 높이', 'mast height', '높이']
        if any(keyword in message_lower for keyword in mast_keywords):
            return self._format_specific_dimension(yacht, 'mastHeight', '마스트 높이 (Mast Height)')
        
        # 세일링/돛 면적
        sail_keywords = ['세일링', 'sailing', '돛', 'sail', '돛 면적', 'sail area', '세일 면적', '넓이', '면적']
        if any(keyword in message_lower for keyword in sail_keywords):
            return self._format_yacht_sail_area(yacht)
        
        # 크기/치수 질문 (전체)
        size_keywords = ['크기', '치수', '수치', 'dimension']
        if any(keyword in message_lower for keyword in size_keywords):
            return self._format_yacht_dimensions(yacht)
        
        # 2. 부품 개수 질문
        parts_count_keywords = ['부품', '부품 개수', '부품 수', 'parts', '몇 개', '개수']
        if any(keyword in message_lower for keyword in parts_count_keywords):
            parts = self._get_yacht_parts(yacht_name)
            parts_count = len(parts) if isinstance(parts, list) else 0
            if parts_count > 0:
                return f"📦 **{yacht_name} 부품 정보**\n\n총 **{parts_count}개**의 부품이 등록되어 있습니다.\n\n더 자세한 정보를 원하시면 '{yacht_name} 분석해줘'라고 물어보세요."
            else:
                return f"📦 **{yacht_name} 부품 정보**\n\n현재 등록된 부품이 없습니다.\n\n부품 정보를 추가하려면 PDF 매뉴얼을 업로드해주세요."
        
        # 3. 엔진 질문
        engine_keywords = ['엔진', 'engine', '모터', 'motor', '동력', '파워']
        if any(keyword in message_lower for keyword in engine_keywords):
            return self._format_yacht_engine_info(yacht)
        
        # 4. 정비/유지보수 질문 ✨ 새로 추가
        maintenance_keywords = ['정비', '유지보수', '관리', '점검', '교체', '주기', 'maintenance', 'repair', 'service', '고장', '수리', '언제']
        if any(keyword in message_lower for keyword in maintenance_keywords):
            return self._format_yacht_maintenance_info(yacht, yacht_name)
        
        # 5. 부품 질문 (특정 부품)
        parts_keywords = ['부품', 'parts', '컴포넌트', 'component']
        if any(keyword in message_lower for keyword in parts_keywords):
            parts = self._get_yacht_parts(yacht_name)
            if isinstance(parts, list) and len(parts) > 0:
                # 부품 목록 반환
                return self._format_yacht_parts_list(yacht_name, parts)
            else:
                return f"📦 **{yacht_name} 부품 정보**\n\n현재 등록된 부품이 없습니다.\n\n부품 정보를 추가하려면 PDF 매뉴얼을 업로드해주세요."
        
        # 6. 제조사 질문
        manufacturer_keywords = ['제조사', 'manufacturer', '만든', '누가']
        if any(keyword in message_lower for keyword in manufacturer_keywords):
            manufacturer = yacht.get('manufacturer', 'N/A')
            return f"🏭 **{yacht_name} 제조사**\n\n제조사: **{manufacturer}**"
        
        # 7. 타입 질문
        type_keywords = ['타입', '유형', '종류', 'type', '어떤']
        if any(keyword in message_lower for keyword in type_keywords):
            yacht_type = yacht.get('type', 'N/A')
            return f"🏷️ **{yacht_name} 유형**\n\n유형: **{yacht_type}**"
        
        # 8. 기본 정보 (간단한 질문)
        info_keywords = ['정보', '스펙', '사양', '알려줘', '뭐야', '어때']
        if any(keyword in message_lower for keyword in info_keywords) and len(user_message.split()) <= 5:
            # 매우 간단한 질문만 처리 (예: "TP52 정보", "Farr 40 알려줘")
            return self._format_basic_yacht_info(yacht)
        
        return None
    
    def _format_basic_yacht_info(self, yacht: Dict) -> str:
        """요트 기본 정보 간단 포맷팅"""
        model_name = yacht.get('name', 'Unknown')
        manufacturer = yacht.get('manufacturer', 'N/A')
        yacht_type = yacht.get('type', 'N/A')
        
        response = f"📋 **{model_name} 기본 정보**\n\n"
        response += f"제조사: {manufacturer}\n"
        response += f"유형: {yacht_type}\n\n"
        
        # 치수 정보가 있으면 간단히 표시
        dim = yacht.get('dimensions', {})
        if dim:
            if dim.get('loa'):
                loa = dim['loa']
                if isinstance(loa, dict):
                    response += f"전장 (LOA): {loa.get('display', loa.get('value', 'N/A'))}\n"
                else:
                    response += f"전장 (LOA): {loa}\n"
            if dim.get('beam'):
                beam = dim['beam']
                if isinstance(beam, dict):
                    response += f"폭 (Beam): {beam.get('display', beam.get('value', 'N/A'))}\n"
                else:
                    response += f"폭 (Beam): {beam}\n"
        
        response += f"\n💡 더 자세한 정보를 원하시면 '{model_name} 분석해줘'라고 물어보세요."
        
        return response
    
    def _format_yacht_engine_info(self, yacht: Dict) -> str:
        """요트 엔진 정보 포맷팅"""
        model_name = yacht.get('name', 'Unknown')
        
        # Schema 5.0: yachtSpecs.standard.engine 경로로 검색
        yacht_specs = yacht.get('yachtSpecs', {})
        standard_specs = yacht_specs.get('standard', {})
        engine = standard_specs.get('engine', {})
        
        # 추가 정보도 확인 (additional에 엔진 정보가 있을 수 있음)
        additional_specs = yacht_specs.get('additional', {})
        
        # 엔진 정보 수집
        engine_type = engine.get('type') or additional_specs.get('engineType') or None
        engine_power = engine.get('power') or additional_specs.get('enginePower') or additional_specs.get('nominalMaximumPropulsionPower') or None
        engine_model = engine.get('model') or additional_specs.get('engineModel') or None
        
        # 정보가 하나도 없으면
        if not engine_type and not engine_power and not engine_model:
            return f"🔧 **{model_name} 엔진 정보**\n\n등록된 엔진 정보가 없습니다.\n\n엔진 정보를 추가하려면 PDF 매뉴얼을 업로드해주세요."
        
        response = f"🔧 **{model_name} 엔진 정보**\n\n"
        
        if engine_type:
            response += f"**타입**: {engine_type}\n"
        if engine_power:
            response += f"**출력**: {engine_power}\n"
        if engine_model:
            response += f"**모델**: {engine_model}\n"
        
        # 추가 엔진 관련 정보가 있으면 표시
        if additional_specs.get('maximumRecommendedEngineSizeWeight'):
            response += f"**권장 엔진 중량**: {additional_specs['maximumRecommendedEngineSizeWeight']}\n"
        
        response += f"\n💡 더 자세한 정보를 원하시면 '{model_name} 분석해줘'라고 물어보세요."
        
        return response
    
    def _format_yacht_maintenance_info(self, yacht: Dict, yacht_name: str) -> str:
        """요트 정비/유지보수 정보 포맷팅 ✨ 새로 추가"""
        model_name = yacht.get('name', 'Unknown')
        
        # 부품 정보에서 정비 주기 추출
        parts = self._get_yacht_parts(yacht_name)
        
        if not parts or len(parts) == 0:
            return f"🔧 **{model_name} 정비 정보**\n\n등록된 부품 및 정비 정보가 없습니다.\n\n정비 정보를 추가하려면 PDF 매뉴얼을 업로드해주세요."
        
        response = f"🔧 **{model_name} 정비 및 유지보수 정보**\n\n"
        
        # 부품별 정비 주기 정리
        maintenance_schedule = {}
        for part in parts:
            if isinstance(part, dict):
                interval = part.get('interval') or part.get('maintenanceInterval')
                if interval:
                    category = part.get('category', '기타')
                    part_name = part.get('name', 'Unknown')
                    
                    if category not in maintenance_schedule:
                        maintenance_schedule[category] = []
                    
                    # interval 숫자 추출
                    interval_value = interval
                    if isinstance(interval, str):
                        import re
                        match = re.search(r'(\d+)', interval)
                        if match:
                            interval_value = int(match.group(1))
                    
                    maintenance_schedule[category].append({
                        'name': part_name,
                        'interval': interval_value,
                        'interval_display': f"{interval}개월" if isinstance(interval, int) else str(interval)
                    })
        
        if not maintenance_schedule:
            response += "정비 주기 정보가 없습니다.\n\n"
            response += f"총 **{len(parts)}개**의 부품이 등록되어 있지만, 정비 주기가 명시되지 않았습니다.\n\n"
        else:
            response += f"**부품별 정비 주기** (총 {len(parts)}개 부품)\n\n"
            
            # 카테고리별로 정리
            for category, items in sorted(maintenance_schedule.items()):
                response += f"**📦 {category}**\n"
                
                # 정비 주기별로 정렬
                sorted_items = sorted(items, key=lambda x: x['interval'] if isinstance(x['interval'], int) else 999)
                
                for item in sorted_items[:5]:  # 각 카테고리당 최대 5개
                    response += f"  • {item['name']}: {item['interval_display']}마다 점검\n"
                
                if len(items) > 5:
                    response += f"  ... 외 {len(items) - 5}개 부품\n"
                
                response += "\n"
        
        # 추가 정보 (maintenance 섹션이 있으면 표시)
        maintenance_info = yacht.get('maintenance', [])
        if maintenance_info and len(maintenance_info) > 0:
            response += "**🔍 추가 정비 정보**\n\n"
            for maint in maintenance_info[:5]:  # 최대 5개
                if isinstance(maint, dict):
                    task = maint.get('task') or maint.get('name', 'Unknown')
                    interval = maint.get('interval', '')
                    method = maint.get('method', '')
                    
                    response += f"**{task}**\n"
                    if interval:
                        response += f"  주기: {interval}\n"
                    if method:
                        response += f"  방법: {method[:100]}\n"  # 100자로 제한
                    response += "\n"
        
        response += "\n💡 정비 관련 궁금한 점은 언제든 물어보세요!"
        response += f"\n📊 전체 부품 목록: '{model_name} 부품'"
        response += f"\n📖 상세 분석: '{model_name} 분석해줘'"
        
        return response
    
    def _format_yacht_parts_list(self, yacht_name: str, parts: List[Dict]) -> str:
        """요트 부품 목록 포맷팅"""
        if not parts or len(parts) == 0:
            return f"📦 **{yacht_name} 부품 정보**\n\n등록된 부품이 없습니다."
        
        response = f"📦 **{yacht_name} 부품 목록**\n\n"
        response += f"총 **{len(parts)}개**의 부품이 등록되어 있습니다.\n\n"
        
        # 카테고리별로 그룹화
        categories = {}
        for part in parts[:20]:  # 최대 20개만 표시
            if isinstance(part, dict):
                category = part.get('category', '기타')
                if category not in categories:
                    categories[category] = []
                categories[category].append(part.get('name', 'Unknown'))
        
        for category, part_names in categories.items():
            response += f"**{category}**: {', '.join(part_names[:5])}"
            if len(part_names) > 5:
                response += f" 외 {len(part_names) - 5}개"
            response += "\n"
        
        if len(parts) > 20:
            response += f"\n... 외 {len(parts) - 20}개 부품 더 있음\n"
        
        response += f"\n💡 전체 부품 목록을 보려면 '{yacht_name} 분석해줘'라고 물어보세요."
        
        return response
    
    def _format_specific_dimension(self, yacht: Dict, dimension_key: str, dimension_name: str) -> str:
        """특정 치수 요소만 포맷팅"""
        model_name = yacht.get('name', 'Unknown')
        
        # Schema 5.0: yachtSpecs.standard.dimensions 경로로 검색
        yacht_specs = yacht.get('yachtSpecs', {})
        standard_specs = yacht_specs.get('standard', {})
        dim = standard_specs.get('dimensions', {})
        
        # 추가 정보도 확인 (detailedDimensions에 더 상세한 정보가 있을 수 있음)
        detailed_dim = yacht.get('detailedDimensions', {})
        
        dimension_data = dim.get(dimension_key) or detailed_dim.get(dimension_key)
        if not dimension_data:
            return f"📏 **{model_name} {dimension_name}**\n\n등록된 {dimension_name} 정보가 없습니다.\n\n정보를 추가하려면 PDF 매뉴얼을 업로드해주세요."
        
        response = f"📏 **{model_name} {dimension_name}**\n\n"
        
        if isinstance(dimension_data, dict):
            value = dimension_data.get('value', '')
            unit = dimension_data.get('unit', '')
            display = dimension_data.get('display', f"{value}{unit}")
            response += f"**{dimension_name}**: {display}\n"
        else:
            response += f"**{dimension_name}**: {dimension_data}\n"
        
        response += f"\n💡 더 자세한 치수 정보를 원하시면 '{model_name} 크기' 또는 '{model_name} 치수'라고 물어보세요."
        
        return response
    
    def _format_yacht_sail_area(self, yacht: Dict) -> str:
        """요트 돛 면적 정보 포맷팅"""
        model_name = yacht.get('name', 'Unknown')
        
        # Schema 5.0: yachtSpecs.standard.sailArea 경로로 검색
        yacht_specs = yacht.get('yachtSpecs', {})
        standard_specs = yacht_specs.get('standard', {})
        sail_area = standard_specs.get('sailArea', {})
        
        # sailInventory도 확인 (더 상세한 정보)
        sail_inventory = yacht.get('sailInventory', {})
        
        if not sail_area and not sail_inventory:
            return f"⛵ **{model_name} 돛 면적**\n\n등록된 돛 면적 정보가 없습니다.\n\n정보를 추가하려면 PDF 매뉴얼을 업로드해주세요."
        
        response = f"⛵ **{model_name} 돛 면적 (Sail Area)**\n\n"
        
        # mainsail (메인 세일)
        main = sail_area.get('mainsail') or sail_area.get('main')
        if main:
            if isinstance(main, dict):
                response += f"**메인 세일 (Mainsail)**: {main.get('value', '')}{main.get('unit', '')}\n"
            else:
                response += f"**메인 세일 (Mainsail)**: {main}\n"
        
        # genoa (제노아)
        genoa = sail_area.get('genoa')
        if genoa:
            if isinstance(genoa, dict):
                response += f"**제노아 (Genoa)**: {genoa.get('value', '')}{genoa.get('unit', '')}\n"
            else:
                response += f"**제노아 (Genoa)**: {genoa}\n"
        
        # jib (지브)
        jib = sail_area.get('jib')
        if jib:
            if isinstance(jib, dict):
                response += f"**지브 (Jib)**: {jib.get('value', '')}{jib.get('unit', '')}\n"
            else:
                response += f"**지브 (Jib)**: {jib}\n"
        
        # spinnaker (스피나커)
        spinnaker = sail_area.get('spinnaker')
        if spinnaker:
            if isinstance(spinnaker, dict):
                response += f"**스피나커 (Spinnaker)**: {spinnaker.get('value', '')}{spinnaker.get('unit', '')}\n"
            else:
                response += f"**스피나커 (Spinnaker)**: {spinnaker}\n"
        
        # total (총 면적)
        total = sail_area.get('total')
        if total:
            if isinstance(total, dict):
                display = total.get('display', f"{total.get('value', '')}{total.get('unit', '')}")
                response += f"**총 면적 (Total)**: {display}\n"
            else:
                response += f"**총 면적 (Total)**: {total}\n"
        
        response += f"\n💡 더 자세한 정보를 원하시면 '{model_name} 분석해줘'라고 물어보세요."
        
        return response
    
    def _format_yacht_dimensions(self, yacht: Dict) -> str:
        """요트 치수 정보 포맷팅"""
        model_name = yacht.get('name', 'Unknown')
        
        # Schema 5.0: yachtSpecs.standard.dimensions 경로로 검색
        yacht_specs = yacht.get('yachtSpecs', {})
        standard_specs = yacht_specs.get('standard', {})
        dim = standard_specs.get('dimensions', {})
        
        # detailedDimensions도 확인
        detailed_dim = yacht.get('detailedDimensions', {})
        
        response = f"📏 **{model_name} 크기 정보**\n\n"
        response += "**기본 치수**\n"
        
        # LOA (전장)
        loa = dim.get('LOA') or detailed_dim.get('LOA')
        if loa:
            if isinstance(loa, dict):
                response += f"- LOA (전장): {loa.get('display', loa.get('value', ''))}\n"
            else:
                response += f"- LOA (전장): {loa}\n"
        
        # LWL (수선장)
        lwl = dim.get('LWL') or dim.get('Lh') or detailed_dim.get('hullLength')
        if lwl:
            if isinstance(lwl, dict):
                response += f"- LWL (수선장): {lwl.get('display', lwl.get('value', ''))}\n"
            else:
                response += f"- LWL (수선장): {lwl}\n"
        
        # Beam (폭)
        beam = dim.get('Beam') or detailed_dim.get('beam')
        if beam:
            if isinstance(beam, dict):
                response += f"- Beam (폭): {beam.get('display', beam.get('value', ''))}\n"
            else:
                response += f"- Beam (폭): {beam}\n"
        
        # Draft (흘수)
        draft = dim.get('Draft') or detailed_dim.get('draughtDeepKeel')
        if draft:
            if isinstance(draft, dict):
                response += f"- Draft (흘수): {draft.get('display', draft.get('value', ''))}\n"
            else:
                response += f"- Draft (흘수): {draft}\n"
        
        # Displacement (배수량)
        disp = dim.get('Displacement') or detailed_dim.get('displacement')
        if disp:
            if isinstance(disp, dict):
                response += f"- Displacement (배수량): {disp.get('display', disp.get('value', ''))}\n"
            else:
                response += f"- Displacement (배수량): {disp}\n"
        
        # Mast Height (마스트 높이)
        mast = dim.get('mastHeight') or detailed_dim.get('airDraftClassicalMast')
        if mast:
            if isinstance(mast, dict):
                response += f"- Mast Height (마스트 높이): {mast.get('display', mast.get('value', ''))}\n"
            else:
                response += f"- Mast Height (마스트 높이): {mast}\n"
        
        return response
    
    def _format_full_yacht_info(self, yacht: Dict) -> str:
        """요트 전체 정보 포맷팅"""
        model_name = yacht.get('name', 'Unknown')
        response = f"🛥️ **{model_name}** - 상세 정보\n\n"
        
        if yacht.get('manufacturer'):
            response += f"제조사: {yacht['manufacturer']}\n"
        if yacht.get('type'):
            response += f"타입: {yacht['type']}\n"
        if yacht.get('designer'):
            response += f"디자이너: {yacht['designer']}\n"
        if yacht.get('year'):
            response += f"제작년도: {yacht['year']}\n"
        
        response += "\n"
        
        # 치수 정보
        dim = yacht.get('dimensions', {})
        if dim:
            response += "📏 **치수**\n"
            if dim.get('loa'):
                loa = dim['loa']
                if isinstance(loa, dict):
                    response += f"- LOA: {loa.get('display', loa.get('value', ''))}\n"
                else:
                    response += f"- LOA: {loa}\n"
            if dim.get('beam'):
                beam = dim['beam']
                if isinstance(beam, dict):
                    response += f"- Beam (폭): {beam.get('display', beam.get('value', ''))}\n"
                else:
                    response += f"- Beam (폭): {beam}\n"
            if dim.get('draft'):
                draft = dim['draft']
                if isinstance(draft, dict):
                    response += f"- Draft (흘수): {draft.get('display', draft.get('value', ''))}\n"
                else:
                    response += f"- Draft (흘수): {draft}\n"
            if dim.get('displacement'):
                disp = dim['displacement']
                if isinstance(disp, dict):
                    response += f"- Displacement (배수량): {disp.get('display', disp.get('value', ''))}\n"
                else:
                    response += f"- Displacement (배수량): {disp}\n"
            if dim.get('mastHeight'):
                mast = dim['mastHeight']
                if isinstance(mast, dict):
                    response += f"- Mast Height: {mast.get('display', mast.get('value', ''))}\n"
                else:
                    response += f"- Mast Height: {mast}\n"
            response += "\n"
        
        # 돛 면적
        sail_area = yacht.get('sailArea', {})
        if sail_area:
            response += "⛵ **돛 면적**\n"
            if sail_area.get('main'):
                main = sail_area['main']
                if isinstance(main, dict):
                    response += f"- Main: {main.get('value', '')} {main.get('unit', '')}\n"
                else:
                    response += f"- Main: {main} m²\n"
            if sail_area.get('jib'):
                jib = sail_area['jib']
                if isinstance(jib, dict):
                    response += f"- Jib: {jib.get('value', '')} {jib.get('unit', '')}\n"
                else:
                    response += f"- Jib: {jib} m²\n"
            if sail_area.get('spinnaker'):
                spin = sail_area['spinnaker']
                if isinstance(spin, dict):
                    response += f"- Spinnaker: {spin.get('value', '')} {spin.get('unit', '')}\n"
                else:
                    response += f"- Spinnaker: {spin} m²\n"
            if sail_area.get('total'):
                total = sail_area['total']
                if isinstance(total, dict):
                    response += f"- Total: {total.get('display', total.get('value', ''))}\n"
                else:
                    response += f"- Total: {total} m²\n"
            response += "\n"
        
        # 엔진 정보
        engine = yacht.get('engine', {})
        if engine:
            response += "🔧 **엔진**\n"
            if engine.get('type'):
                response += f"- Type: {engine['type']}\n"
            if engine.get('power'):
                response += f"- Power: {engine['power']}\n"
            if engine.get('model'):
                response += f"- Model: {engine['model']}\n"
        
        return response
    
    def _build_context(self) -> str:
        """대화 컨텍스트 구성"""
        context = self.system_prompt + "\n\n**대화 기록:**\n"
        recent_history = self.chat_history[-10:]
        
        for msg in recent_history:
            role = "사용자" if msg["role"] == "user" else "어시스턴트"
            context += f"\n{role}: {msg['content']}\n"
        
        return context
    
    def _handle_file_upload(self, file_path: str) -> str:
        """파일 업로드 및 분석 처리 (PDF, Word, HWP, Excel, PPTX 등)"""
        try:
            file_name = os.path.basename(file_path)
            file_ext = self._get_file_extension(file_path)
            
            # 지원되는 파일 형식 확인
            if not self._is_supported_file(file_path):
                return f"❌ 지원되지 않는 파일 형식입니다.\n\n지원 형식: PDF, Word (.docx, .doc), HWP, 텍스트 (.txt), Excel (.xlsx, .xls), PowerPoint (.pptx, .ppt)"
            
            print(f"📄 파일 분석 시작: {file_name} ({file_ext})", flush=True)
            
            # 파일 분석 시작 메시지
            analyzing_msg = f"📄 {file_name} 문서를 분석 중입니다...\n잠시만 기다려주세요! ⏳"
            
            self.chat_history.append({
                "role": "user",
                "content": f"[파일 업로드: {file_name}]",
                "timestamp": datetime.now().isoformat()
            })
            
            self.chat_history.append({
                "role": "assistant",
                "content": analyzing_msg,
                "timestamp": datetime.now().isoformat()
            })
            
            # 파일 형식에 따라 텍스트 추출
            print("📖 텍스트 추출 중...", flush=True)
            extracted_text = self._extract_text_from_file(file_path)
            
            if not extracted_text or len(extracted_text.strip()) < 100:
                return f"❌ {file_name}에서 텍스트를 추출할 수 없습니다.\n\n파일이 손상되었거나 암호화되어 있을 수 있습니다."
            
            # 분석 실행
            print(f"✅ 텍스트 추출 완료 ({len(extracted_text)} 문자)", flush=True)
            print("🤖 AI 분석 시작...", flush=True)
            analysis_result = self._analyze_document_directly(file_path, extracted_text)
            
            # 분석 결과 확인
            if "error" in analysis_result:
                error_msg = f"❌ 문서 분석 중 오류가 발생했습니다:\n{analysis_result.get('error', '알 수 없는 오류')}"
                self.chat_history.append({
                    "role": "assistant",
                    "content": error_msg,
                    "timestamp": datetime.now().isoformat()
                })
                return error_msg
            
            # 분석 결과를 요트 등록 형식으로 변환
            registration_data = self._convert_analysis_to_registration(analysis_result)
            
            # 등록 완료 메시지 생성
            completion_msg = self._generate_registration_completion_message(analysis_result, registration_data)
            
            # 등록 데이터 저장 (메모리 + JSON 파일)
            self.current_yacht_registration = registration_data
            
            # JSON 파일로 저장
            self._save_registration_to_json(registration_data, analysis_result)
            
            # 대화 히스토리에 추가
            self.chat_history.append({
                "role": "assistant",
                "content": completion_msg,
                "timestamp": datetime.now().isoformat()
            })
            
            print(f"✅ {file_name} 분석 및 등록 준비 완료!")
            
            return completion_msg
            
        except Exception as e:
            error_msg = f"❌ 파일 처리 중 오류가 발생했습니다: {str(e)}"
            print(f"❌ Error: {e}")
            import traceback
            traceback.print_exc()
            
            self.chat_history.append({
                "role": "assistant",
                "content": error_msg,
                "timestamp": datetime.now().isoformat()
            })
            
            return error_msg
    
    def _analyze_document_directly(self, file_path: str, extracted_text: str) -> Dict:
        """문서 직접 분석 (yacht_document_analyzer 없이)"""
        # 텍스트가 너무 길면 앞부분만 사용
        if len(extracted_text) > 30000:
            extracted_text = extracted_text[:30000] + "\n\n[... 텍스트가 너무 길어 일부만 분석합니다 ...]"
        
        if not self.has_gemini:
            return {
                "error": "문서 분석 기능은 Gemini API가 필요합니다.",
                "fileInfo": {
                    "fileName": os.path.basename(file_path),
                    "filePath": file_path
                }
            }
        
        # 분석 프롬프트 (완전한 버전 5.0)
        prompt = f"""다음은 요트 매뉴얼 또는 부품 정보 문서에서 추출한 텍스트입니다:

{extracted_text}

---

## 📋 작업 지시사항 (Schema Version 5.0)

매뉴얼에서 발견한 **모든 정보**를 최대한 상세하게 추출하세요.

---

### ✅ 섹션 1: 문서 기본 정보
```json
"documentInfo": {{
  "title": "문서 제목",
  "yachtModel": "요트 모델명",
  "manufacturer": "제조사",
  "documentType": "Owner's Manual / Parts List / Technical Specifications / Class Rules"
}}
```

---

### ✅ 섹션 2: 요트 기본 스펙
**standard (표준 필드):**
- dimensions: LOA, LWL, Beam, Draft, Displacement, mastHeight
- engine: type, power, model
- sailArea: mainsail, jib, spinnaker, total

**additional (발견한 모든 추가 정보):**
- 위 standard에 없는 모든 스펙을 키-값으로 저장
- 키 이름: camelCase (예: keelWeight, fuelCapacity)
- 신뢰도: _confidence_{{키이름}}: "high" / "medium" / "low"

---

### ✅ 섹션 3: 상세 치수 (Detailed Dimensions)
**모든 치수 정보를 추출하세요:**
- LWL, BOA, freeboard (bow/midship/stern)
- headroom (saloon/cabins/galley)
- ballastWeight, ballastRatio
- keel dimensions, rudder dimensions
- boom length, pole length
- 신뢰도: _confidence_{{키이름}}

---

### ✅ 섹션 4: 외관 (Exterior)

**🔑 중요: 모든 항목에 고유 ID 부여!**

**ID 생성 규칙:**
- Hull: `ext-hull-01`
- Keel: `ext-hull-keel-01`
- Rudder: `ext-hull-rudder-01`
- Deck: `ext-deck-01`
- Cockpit: `ext-deck-cockpit-01`
- Windows: `ext-window-{{location}}-{{number}}`
- Hatches: `ext-hatch-{{location}}-{{number}}`

**hull:**
```json
{{
  "id": "ext-hull-01",
  "name": "Hull",
  "category": "Structure",
  "manufacturer": "...",
  "specifications": {{
    "type": "Monohull",
    "material": "GRP / Fiberglass / Carbon",
    "color": "...",
    "thickness": "...",
    "gelcoatType": "...",
    "coreType": "Balsa / Foam / Solid",
    "_confidence_material": "high",
    "_additional": {{}}
  }},
  "subComponents": [
    {{
      "id": "ext-hull-keel-01",
      "parentId": "ext-hull-01",
      "name": "Keel",
      "category": "Hull Structure",
      "specifications": {{
        "type": "Fin / Bulb / Canting",
        "material": "Lead / Iron / Composite",
        "weight": "...",
        "draft": "...",
        "attachmentMethod": "..."
      }},
      "maintenanceDetails": {{
        "interval": 12,
        "inspectionItems": ["Keel bolts", "Corrosion", "Leakage"],
        "commonIssues": "...",
        "repairCost": "..."
      }}
    }},
    {{
      "id": "ext-hull-rudder-01",
      "parentId": "ext-hull-01",
      "name": "Rudder",
      "specifications": {{
        "type": "Spade / Skeg-mounted",
        "material": "...",
        "dimensions": "..."
      }}
    }}
  ]
}}
```

**deck, windows, hatches:** 동일한 구조로 추출

---

### ✅ 섹션 5: 앵커 시스템 (Ground Tackle)

**ID 생성 규칙:**
- Anchors: `anchor-{{type}}-{{number}}`
- Chain: `anchor-chain-01`
- Windlass: `anchor-windlass-01`
- Windlass parts: `anchor-windlass-{{part}}-{{number}}`

```json
{{
  "anchors": [
    {{
      "id": "anchor-primary-01",
      "name": "Primary Anchor",
      "type": "Delta / CQR / Fortress / Rocna",
      "manufacturer": "...",
      "model": "...",
      "specifications": {{
        "weight": "... kg",
        "material": "...",
        "holdingPower": "... kg",
        "_confidence_weight": "high"
      }}
    }}
  ],
  "chain": {{
    "id": "anchor-chain-01",
    "name": "Anchor Chain",
    "specifications": {{
      "material": "Galvanized steel / Stainless",
      "diameter": "... mm",
      "length": "... m",
      "grade": "..."
    }}
  }},
  "windlass": {{
    "id": "anchor-windlass-01",
    "name": "Windlass",
    "manufacturer": "...",
    "specifications": {{
      "type": "Electric / Manual / Hydraulic",
      "power": "...",
      "maxPull": "..."
    }},
    "subComponents": [...]
  }}
}}
```

---

### ✅ 섹션 6: 돛 목록 (Sail Inventory)

**ID 생성 규칙:**
- Mainsail: `sail-main-01`
- Genoa: `sail-genoa-{{size}}-01`
- Spinnaker: `sail-spinnaker-01`

```json
[
  {{
    "id": "sail-main-01",
    "name": "Mainsail",
    "category": "Sails",
    "manufacturer": "North Sails / Quantum / UK Sailmakers",
    "model": "...",
    "specifications": {{
      "area": "... m²",
      "luffLength": "... m",
      "footLength": "... m",
      "material": "Dacron / Mylar / 3Di / Carbon",
      "weight": "... kg",
      "year": "...",
      "reefingPoints": 2,
      "numberOfBattens": 4,
      "condition": "Excellent / Good / Fair / Poor",
      "_confidence_area": "high"
    }},
    "subComponents": [
      {{
        "id": "sail-main-slides-01",
        "parentId": "sail-main-01",
        "name": "Sail Slides",
        "specifications": {{
          "type": "...",
          "quantity": 12
        }}
      }}
    ],
    "maintenanceDetails": {{
      "interval": 6,
      "inspectionItems": ["Stitching", "UV cover", "Battens"],
      "repairCost": "..."
    }}
  }}
]
```

---

### ✅ 섹션 7: 갑판 장비 (Deck Equipment)

**ID 생성 규칙:**
- Winches: `deck-winch-{{location}}-{{number}}`
- Cleats: `deck-cleat-{{location}}-{{number}}`
- Blocks: `deck-block-{{type}}-{{number}}`

```json
{{
  "winches": [
    {{
      "id": "deck-winch-primary-port-01",
      "name": "Primary Winch Port",
      "manufacturer": "Harken / Lewmar / Andersen",
      "model": "...",
      "category": "Deck Hardware",
      "specifications": {{
        "location": "Cockpit coaming port",
        "type": "Two-speed self-tailing",
        "gearRatio": "...:1",
        "drumDiameter": "... mm",
        "maxLoad": "... kg",
        "weight": "... kg",
        "material": "Aluminum / Bronze",
        "partNumber": "..."
      }},
      "subComponents": [
        {{
          "id": "deck-winch-primary-port-handle-01",
          "parentId": "deck-winch-primary-port-01",
          "name": "Winch Handle",
          "specifications": {{"length": "... mm"}}
        }}
      ],
      "maintenanceDetails": {{
        "interval": 12,
        "inspectionItems": ["Pawls", "Gears", "Drum"],
        "lubricationType": "Marine winch grease",
        "repairCost": "$50-200"
      }}
    }}
  ],
  "cleats": [...],
  "blocks": [...],
  "stanchions": {{...}},
  "steeringSystem": {{...}}
}}
```

---

### ✅ 섹션 8: 시설물 (Accommodations)

**ID 생성 규칙:**
- Galley: `accom-galley-01`
- Galley components: `accom-galley-{{component}}-01`
- Cabins: `accom-cabin-{{location}}-01`
- Heads: `accom-head-{{location}}-01`

```json
{{
  "summary": {{
    "cabins": 3,
    "berths": 6,
    "heads": 2,
    "showers": 1
  }},
  "galley": {{
    "id": "accom-galley-01",
    "name": "Galley",
    "location": "Port / Starboard / Center",
    "specifications": {{
      "dimensions": "... x ... m",
      "counterMaterial": "Corian / Laminate",
      "storageVolume": "... L"
    }},
    "components": [
      {{
        "id": "accom-galley-stove-01",
        "parentId": "accom-galley-01",
        "name": "Stove",
        "manufacturer": "Force 10 / Eno / Dometic",
        "model": "...",
        "specifications": {{
          "type": "2-burner / 3-burner gas / electric",
          "fuelType": "LPG / CNG",
          "power": "...",
          "gimbalMount": true
        }},
        "maintenanceDetails": {{...}}
      }},
      {{
        "id": "accom-galley-fridge-01",
        "parentId": "accom-galley-01",
        "name": "Refrigerator",
        "manufacturer": "Isotherm / Frigoboat",
        "specifications": {{
          "capacity": "... L",
          "type": "12V compressor / Eutectic",
          "powerConsumption": "... A"
        }}
      }},
      {{
        "id": "accom-galley-sink-01",
        "name": "Galley Sink",
        "specifications": {{
          "material": "Stainless steel",
          "numberOfBowls": 1
        }}
      }}
    ]
  }},
  "cabins": [
    {{
      "id": "accom-cabin-master-01",
      "name": "Master Cabin",
      "location": "Aft / Forward",
      "specifications": {{
        "berthSize": "Queen / Double / Twin",
        "headroom": "... m",
        "privateHead": true
      }},
      "components": [...]
    }}
  ],
  "heads": [
    {{
      "id": "accom-head-forward-01",
      "name": "Forward Head",
      "specifications": {{
        "shower": true,
        "showerType": "Wet head / Separate"
      }},
      "components": [
        {{
          "id": "accom-head-forward-toilet-01",
          "name": "Marine Toilet",
          "manufacturer": "Jabsco / Raritan / Tecma",
          "model": "...",
          "specifications": {{
            "type": "Manual / Electric",
            "discharge": "Overboard / Holding tank"
          }},
          "maintenanceDetails": {{...}}
        }}
      ]
    }}
  ]
}}
```

---

### ✅ 섹션 9: 수조 (Tanks)

**ID: `tank-{{type}}-{{number}}`**

```json
{{
  "fuel": {{
    "id": "tank-fuel-01",
    "name": "Fuel Tank",
    "specifications": {{
      "capacity": "... L",
      "material": "Stainless steel / Aluminum / Plastic",
      "location": "...",
      "fuelType": "Diesel / Gasoline"
    }},
    "subComponents": [...]
  }},
  "freshWater": {{
    "id": "tank-water-01",
    "specifications": {{
      "totalCapacity": "... L",
      "material": "Food-grade polyethylene",
      "numberOfTanks": 2
    }}
  }},
  "holdingTank": {{
    "id": "tank-holding-01",
    "specifications": {{
      "capacity": "... L",
      "pumpout": true
    }}
  }}
}}
```

---

### ✅ 섹션 10: 전기 시스템 (Electrical System)

**ID: `elec-{{category}}-{{component}}-{{number}}`**

```json
{{
  "batteries": {{
    "house": {{
      "id": "elec-battery-house-01",
      "name": "House Battery Bank",
      "manufacturer": "Victron / Lifeline / Trojan",
      "model": "...",
      "specifications": {{
        "type": "AGM / Gel / Lithium / Flooded Lead-Acid",
        "totalCapacity": "... Ah",
        "voltage": "12V / 24V",
        "numberOfBatteries": 2,
        "configuration": "Parallel / Series"
      }}
    }},
    "starter": {{...}}
  }},
  "chargers": [...],
  "solarPanels": {{
    "id": "elec-solar-array-01",
    "specifications": {{
      "totalCapacity": "... W",
      "numberOfPanels": 2
    }},
    "subComponents": [
      {{
        "id": "elec-solar-controller-01",
        "name": "Solar Charge Controller",
        "manufacturer": "Victron / Morningstar",
        "specifications": {{
          "type": "MPPT / PWM",
          "maxPVVoltage": "... V",
          "maxChargeCurrent": "... A"
        }}
      }}
    ]
  }},
  "inverter": {{...}},
  "shoreConnection": {{...}}
}}
```

---

### ✅ 섹션 11: 전자 장비 (Electronics)

**ID: `electron-{{category}}-{{component}}-{{number}}`**

```json
{{
  "navigation": [
    {{
      "id": "electron-nav-chartplotter-01",
      "name": "Chartplotter",
      "manufacturer": "Raymarine / Garmin / Simrad / B&G",
      "model": "...",
      "specifications": {{
        "type": "Multifunction display",
        "screenSize": "... inch",
        "resolution": "...",
        "touchscreen": true,
        "cartography": "Navionics / C-MAP"
      }}
    }},
    {{
      "id": "electron-nav-radar-01",
      "name": "Radar",
      "specifications": {{
        "type": "Doppler / Pulse",
        "range": "... NM"
      }}
    }}
  ],
  "communication": [
    {{
      "id": "electron-comm-vhf-01",
      "name": "VHF Radio",
      "manufacturer": "Standard Horizon / Icom",
      "specifications": {{
        "type": "Fixed mount / Handheld",
        "dsc": true,
        "power": "... W"
      }}
    }}
  ],
  "instruments": [...],
  "autopilot": {{
    "id": "electron-autopilot-01",
    "specifications": {{
      "type": "Hydraulic / Electric / Wind vane",
      "manufacturer": "..."
    }},
    "subComponents": [...]
  }}
}}
```

---

### ✅ 섹션 12: 배관 시스템 (Plumbing System)

**ID: `plumb-{{category}}-{{component}}-{{number}}`**

```json
{{
  "waterMaker": {{
    "id": "plumb-watermaker-01",
    "specifications": {{
      "type": "Reverse osmosis",
      "capacity": "... L/hour"
    }},
    "subComponents": [...]
  }},
  "pumps": [
    {{
      "id": "plumb-pump-freshwater-01",
      "name": "Freshwater Pressure Pump",
      "manufacturer": "Jabsco / Shurflo / Whale",
      "specifications": {{
        "flow": "... L/min",
        "pressure": "... bar"
      }}
    }}
  ],
  "bilgePumps": [
    {{
      "id": "plumb-bilge-primary-01",
      "name": "Primary Bilge Pump",
      "specifications": {{
        "type": "Automatic / Manual",
        "capacity": "... GPH"
      }}
    }}
  ],
  "seacocks": {{
    "id": "plumb-seacocks-01",
    "specifications": {{
      "totalQuantity": 8,
      "material": "Bronze / Marelon"
    }},
    "components": [...]
  }}
}}
```

---

### ✅ 섹션 13: 부품 (Parts) - 통합 리스트

**ID: `part-{{category}}-{{name}}-{{number}}`**

⚠️ **중요: 대체 가능한 부품 (Alternative Parts) 통합 규칙**

**1. 같은 이름의 부품이 여러 제조사로 나열된 경우:**
   - 하나의 부품으로 통합
   - manufacturer 필드에 모든 제조사를 슬래시(/)로 구분
   - model 필드에 모든 모델을 슬래시(/)로 구분

**2. "OR", "alternatively", "/" 키워드 발견 시:**
   - 대체 부품으로 인식하여 하나로 통합
   
**3. 같은 카테고리 + 같은 위치의 부품:**
   - 기능이 같다면 하나로 통합

**올바른 예시:**
```json
{{
  "name": "Primary Winch",
  "manufacturer": "Harken / Lewmar",
  "model": "B480TCR / Ocean Racing 440 / 44",
  "category": "Deck Equipment"
}}
```

**잘못된 예시 (중복 - 하지 마세요):**
```json
[
  {{"name": "Primary Winch", "manufacturer": "Harken", "model": "B480TCR"}},
  {{"name": "Primary Winch", "manufacturer": "Lewmar", "model": "Ocean Racing 440"}},
  {{"name": "Primary Winch", "manufacturer": "Lewmar", "model": "44"}}
]
```

**부품 추출 형식:**
```json
[
  {{
    "id": "part-rigging-mast-01",
    "name": "Mast",
    "manufacturer": "Selden / Z-Spars / Hall Spars",
    "model": "...",
    "interval": 12,
    "category": "Rigging",
    "specifications": {{
      "material": "Aluminum / Carbon",
      "length": "... m",
      "weight": "... kg",
      "partNumber": "...",
      "_confidence_length": "high"
    }},
    "subParts": [...],
    "maintenanceDetails": {{
      "interval": 12,
      "inspectionItems": ["Corrosion", "Bolts", "Wiring"],
      "repairCost": "..."
    }}
  }}
]
```

---

### ✅ 섹션 14: 유지보수 (Maintenance)

```json
[
  {{
    "item": "...",
    "interval": "... 개월",
    "method": "..."
  }}
]
```

---

### ✅ 섹션 15: 분석 결과 (Analysis Result)

```json
{{
  "canExtractText": true/false,
  "canAnalyze": true/false,
  "reason": "요트 매뉴얼이 아닌 경우 이유 설명"
}}
```

---

## 🎯 최종 응답 형식

```json
{{
  "schemaVersion": "5.0",
  "analyzedAt": "2025-11-20T10:30:00Z",
  "documentInfo": {{...}},
  "yachtSpecs": {{
    "standard": {{...}},
    "additional": {{...}}
  }},
  "detailedDimensions": {{...}},
  "exterior": {{...}},
  "groundTackle": {{...}},
  "sailInventory": [...],
  "deckEquipment": {{...}},
  "accommodations": {{...}},
  "tanks": {{...}},
  "electricalSystem": {{...}},
  "electronics": {{...}},
  "plumbingSystem": {{...}},
  "parts": [...],
  "maintenance": [...],
  "analysisResult": {{...}}
}}
```

---

## ⚠️ 중요 규칙

1. **ID 필수**: 모든 항목에 고유 ID 부여
2. **부모-자식 관계**: subComponents/subParts에 parentId 추가
3. **신뢰도**: 중요 필드에 _confidence 추가
4. **확장성**: _additional 필드 활용
5. **중복 방지**: standard에 있는 정보는 additional에 추가 금지
6. **추측 금지**: 불확실하면 null
7. **JSON만**: 다른 설명 불필요

**JSON 형식으로만 응답해주세요.**"""
        
        # Gemini API 호출
        print("🤖 AI 분석 중...")
        response = self.model.generate_content(prompt)
        
        # 응답 파싱
        result_text = response.text
        
        # JSON 추출 (마크다운 코드 블록 제거)
        if "```json" in result_text:
            json_start = result_text.find("```json") + 7
            json_end = result_text.find("```", json_start)
            result_text = result_text[json_start:json_end].strip()
        elif "```" in result_text:
            json_start = result_text.find("```") + 3
            json_end = result_text.find("```", json_start)
            result_text = result_text[json_start:json_end].strip()
        
        # JSON 파싱
        try:
            result = json.loads(result_text)
        except json.JSONDecodeError:
            result = {
                "rawResponse": result_text,
                "error": "JSON 파싱 실패",
                "extractedTextLength": len(extracted_text)
            }
        
        # 파일 정보 추가
        result["fileInfo"] = {
            "fileName": os.path.basename(file_path),
            "filePath": file_path,
            "fileSize": os.path.getsize(file_path)
        }
        
        print("✅ 분석 완료!")
        return result
    
    def _extract_text_from_pdf(self, pdf_path: str) -> str:
        """PDF에서 텍스트 추출 (일반 방법 실패 시 OCR 시도)"""
        text = ""
        
        # 방법 1: PyPDF2로 텍스트 추출 시도
        if HAS_PYPDF2:
            try:
                with open(pdf_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page in pdf_reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                # 텍스트가 충분히 추출되었으면 반환
                if len(text.strip()) > 100:
                    return text
            except Exception as e:
                print(f"⚠️ PyPDF2로 텍스트 추출 실패: {e}")
        
        # 방법 2: pdfplumber로 텍스트 추출 시도
        if HAS_PDFPLUMBER:
            try:
                import pdfplumber
                with pdfplumber.open(pdf_path) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                # 텍스트가 충분히 추출되었으면 반환
                if len(text.strip()) > 100:
                    return text
            except Exception as e:
                print(f"⚠️ pdfplumber로 텍스트 추출 실패: {e}")
        
        # 방법 3: EasyOCR 사용 (스캔된 이미지 PDF인 경우)
        if len(text.strip()) < 100:
            try:
                print("📷 텍스트 추출 실패. OCR을 시도합니다...")
                text = self._extract_text_with_easyocr(pdf_path)
                if len(text.strip()) > 100:
                    print("✅ OCR로 텍스트 추출 성공!")
                    return text
            except Exception as e:
                print(f"⚠️ OCR 실패: {e}")
                print("💡 OCR 패키지를 설치하려면:")
                print("   python install_ocr_local.py")
        
        return text
    
    def _extract_text_with_easyocr(self, pdf_path: str) -> str:
        """EasyOCR을 사용한 텍스트 추출 (스캔된 이미지 PDF용)"""
        try:
            import fitz  # PyMuPDF
            import easyocr
            import numpy as np
            from PIL import Image
            import io
            
            # EasyOCR 초기화 (영어 + 한국어)
            print("   🤖 EasyOCR 초기화 중...")
            reader = easyocr.Reader(['en', 'ko'], gpu=False)
            
            # PDF 열기
            doc = fitz.open(pdf_path)
            total_pages = len(doc)
            print(f"   📄 총 {total_pages}페이지를 OCR 처리 중...")
            
            text = ""
            
            for page_num in range(total_pages):
                # PDF 페이지를 이미지로 변환
                page = doc[page_num]
                pix = page.get_pixmap(dpi=300)
                img_data = pix.tobytes("png")
                img = Image.open(io.BytesIO(img_data))
                
                # numpy 배열로 변환
                img_array = np.array(img)
                
                # OCR 실행
                results = reader.readtext(img_array)
                
                # 결과 텍스트 추출
                page_text = "\n".join([text_result[1] for text_result in results])
                text += f"\n--- Page {page_num + 1} ---\n{page_text}\n"
                
                if (page_num + 1) % 5 == 0:
                    print(f"   진행 중: {page_num + 1}/{total_pages} 페이지")
            
            doc.close()
            print(f"   ✅ OCR 완료: {len(text)} 문자 추출")
            return text
            
        except ImportError as e:
            print(f"❌ 필요한 패키지가 설치되지 않았습니다: {e}")
            print("💡 실행: python install_ocr_local.py")
            return ""
        except Exception as e:
            print(f"❌ OCR 오류: {e}")
            import traceback
            traceback.print_exc()
            return ""
    
    def _extract_text_from_file(self, file_path: str) -> str:
        """파일에서 텍스트 추출 (PDF, Word, HWP, Excel, PPTX 등)"""
        file_ext = self._get_file_extension(file_path)
        
        if file_ext == '.pdf':
            return self._extract_text_from_pdf(file_path)
        elif file_ext in ['.docx', '.doc']:
            return self._extract_text_from_word(file_path)
        elif file_ext == '.hwp':
            return self._extract_text_from_hwp(file_path)
        elif file_ext == '.txt':
            return self._extract_text_from_txt(file_path)
        elif file_ext in ['.xlsx', '.xls']:
            return self._extract_text_from_excel(file_path)
        elif file_ext in ['.pptx', '.ppt']:
            return self._extract_text_from_pptx(file_path)
        else:
            return ""
    
    def _extract_text_from_word(self, file_path: str) -> str:
        """Word 문서에서 텍스트 추출"""
        if not HAS_DOCX:
            print("⚠️ python-docx가 설치되지 않았습니다. pip install python-docx")
            return ""
        
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # 테이블에서도 텍스트 추출
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\n"
            
            return text
        except Exception as e:
            print(f"⚠️ Word 문서 텍스트 추출 실패: {e}")
            return ""
    
    def _extract_text_from_hwp(self, file_path: str) -> str:
        """HWP 파일에서 텍스트 추출"""
        if not HAS_OLEFILE:
            print("⚠️ olefile이 설치되지 않았습니다. pip install olefile")
            return ""
        
        try:
            # HWP 파일은 OLE 형식
            if not olefile.isOleFile(file_path):
                return ""
            
            ole = olefile.OleFileIO(file_path)
            text = ""
            
            # HWP 파일 구조에서 텍스트 추출 시도
            try:
                # Section0 스트림에서 텍스트 추출 시도
                if ole.exists('Section0'):
                    stream = ole.openstream('Section0')
                    data = stream.read()
                    # 한글 인코딩 시도
                    try:
                        text = data.decode('utf-8', errors='ignore')
                    except:
                        try:
                            text = data.decode('cp949', errors='ignore')
                        except:
                            text = data.decode('latin-1', errors='ignore')
            except Exception as e:
                print(f"⚠️ HWP 텍스트 추출 시도 중 오류: {e}")
            
            ole.close()
            
            # 텍스트가 너무 짧으면 실패로 간주
            if len(text.strip()) < 50:
                return ""
            
            return text
        except Exception as e:
            print(f"⚠️ HWP 파일 텍스트 추출 실패: {e}")
            print("💡 HWP 파일은 복잡한 형식이므로 완벽한 추출이 어려울 수 있습니다.")
            return ""
    
    def _extract_text_from_txt(self, file_path: str) -> str:
        """텍스트 파일에서 내용 읽기"""
        try:
            # 여러 인코딩 시도
            encodings = ['utf-8', 'cp949', 'euc-kr', 'latin-1']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            
            # 모두 실패하면 바이너리로 읽기
            with open(file_path, 'rb') as f:
                return f.read().decode('utf-8', errors='ignore')
        except Exception as e:
            print(f"⚠️ 텍스트 파일 읽기 실패: {e}")
            return ""
    
    def _extract_text_from_excel(self, file_path: str) -> str:
        """Excel 파일에서 텍스트 추출"""
        if not HAS_OPENPYXL:
            print("⚠️ openpyxl이 설치되지 않았습니다. pip install openpyxl")
            return ""
        
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            text = ""
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text += f"\n[시트: {sheet_name}]\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
            
            wb.close()
            return text
        except Exception as e:
            print(f"⚠️ Excel 파일 텍스트 추출 실패: {e}")
            return ""
    
    def _extract_text_from_pptx(self, file_path: str) -> str:
        """PowerPoint 파일에서 텍스트 추출"""
        if not HAS_PPTX:
            print("⚠️ python-pptx가 설치되지 않았습니다. pip install python-pptx")
            return ""
        
        try:
            prs = Presentation(file_path)
            text = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\n--- 슬라이드 {slide_num} ---\n"
                
                # 슬라이드의 모든 도형에서 텍스트 추출
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        shape_text = shape.text.strip()
                        if shape_text:
                            text += shape_text + "\n"
                    
                    # 테이블이 있는 경우
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join([cell.text.strip() if cell.text else "" for cell in row.cells])
                            if row_text.strip():
                                text += row_text + "\n"
            
            return text
        except Exception as e:
            print(f"⚠️ PowerPoint 파일 텍스트 추출 실패: {e}")
            return ""
    
    def _list_yachts(self) -> str:
        """요트 목록 반환"""
        yachts = self.yacht_data.get('yachts', [])
        if not yachts:
            return "요트 데이터를 찾을 수 없습니다."
        
        result = f"📋 총 {len(yachts)}개의 요트 모델:\n\n"
        for i, yacht in enumerate(yachts, 1):
            name = yacht.get('name', 'Unknown')
            yacht_type = yacht.get('type', '')
            result += f"  {i}. {name}"
            if yacht_type:
                result += f" ({yacht_type})"
            result += "\n"
        
        return result
    
    def _get_data_info(self) -> str:
        """데이터 정보 반환"""
        yachts = self.yacht_data.get('yachts', [])
        result = "📊 요트 데이터 정보\n"
        result += "=" * 50 + "\n"
        result += f"총 요트 개수: {len(yachts)}개\n"
        result += f"데이터 버전: {self.yacht_data.get('version', 'N/A')}\n"
        result += f"마지막 업데이트: {self.yacht_data.get('lastUpdated', 'N/A')}\n"
        return result
    
    def _get_help(self) -> str:
        """도움말 반환"""
        help_text = """📖 HooAah Yacht 챗봇 도움말

**사용 가능한 명령어:**
- `/list` 또는 `/목록` - 요트 목록 보기
- `/info` 또는 `/정보` - 데이터 정보 보기
- `/help` 또는 `/도움말` - 이 도움말 보기

**질문 예시:**
- "Farr 40 크기 알려줘"
- "레이싱에 좋은 요트 추천해줘"
- "PDF 파일 경로" (PDF 업로드 및 분석)

**PDF 업로드:**
PDF 파일 경로를 입력하면 자동으로 분석합니다.
예: "C:\\Users\\user\\Documents\\manual.pdf"
"""
        return help_text
    
    def _suggest_pdf_upload(self) -> str:
        """PDF 업로드 안내 메시지 (Gemini AI 사용 가능 시)"""
        message = """📄 요트 문서를 등록하세요!

요트 매뉴얼 PDF 파일을 업로드해주시면:
1. 📋 문서를 자동으로 분석합니다
2. ⛵ 요트 정보를 추출합니다
3. 🔧 부품 정보를 정리합니다
4. ✅ 데이터베이스에 등록합니다

**PDF 파일 업로드 방법:**
PDF 파일의 전체 경로를 입력해주세요.

예시:
- `C:\\Users\\user\\Documents\\Sun Odyssey 380 Owners manual.pdf`
- `"C:\\Users\\user\\Documents\\manual.pdf"` (공백이 있는 경우 따옴표 사용)

PDF 파일 경로를 입력해주세요! 📎"""
        
        return message
    
    def _suggest_pdf_upload_without_ai(self) -> str:
        """PDF 업로드 안내 메시지 (Gemini AI 없을 때)"""
        message = """📄 요트 문서를 등록하세요!

요트 매뉴얼 PDF 파일을 업로드해주시면 자동으로 분석하고 등록합니다.

**PDF 파일 업로드 방법:**
PDF 파일의 전체 경로를 입력해주세요.

예시:
- `C:\\Users\\user\\Documents\\Sun Odyssey 380 Owners manual.pdf`
- `"C:\\Users\\user\\Documents\\manual.pdf"` (공백이 있는 경우 따옴표 사용)

⚠️ **참고**: PDF 분석 기능은 Gemini API가 필요합니다.
현재 Gemini API가 설정되지 않아 PDF 분석 기능을 사용할 수 없습니다.

Gemini API 키를 설정하려면:
1. 환경변수에 `GEMINI_API_KEY` 설정
2. 또는 실행 시 `--api-key YOUR_API_KEY` 옵션 사용

PDF 파일 경로를 입력해주세요! 📎"""
        
        return message
    
    def _extract_section_keyword(self, message: str) -> Optional[str]:
        """메시지에서 섹션 키워드 추출
        
        Args:
            message: 사용자 메시지
            
        Returns:
            섹션 키워드 (예: "사용 목적에 따른 적합성 평가") 또는 None
        """
        message_lower = message.lower()
        
        # 섹션 키워드 매핑 (사용자 입력 → 정확한 섹션 제목)
        section_keywords = {
            "사용 목적": "사용 목적에 따른 적합성 평가",
            "적합성": "사용 목적에 따른 적합성 평가",
            "적합성 평가": "사용 목적에 따른 적합성 평가",
            "정비": "관리 및 정비 권장사항",
            "정비 권장": "관리 및 정비 권장사항",
            "관리": "관리 및 정비 권장사항",
            "권장사항": "관리 및 정비 권장사항",
            "특징": "요트의 주요 특징 및 스펙 요약",
            "스펙": "요트의 주요 특징 및 스펙 요약",
            "치수": "치수 및 성능 분석",
            "성능": "치수 및 성능 분석",
            "부품": "부품 구성 및 정비 주기 분석",
            "부품 구성": "부품 구성 및 정비 주기 분석"
        }
        
        # 키워드 검색
        for keyword, section_title in section_keywords.items():
            if keyword in message_lower:
                return section_title
        
        return None
    
    def _extract_yacht_name_from_message(self, message: str) -> Optional[str]:
        """메시지에서 요트 이름 추출 (하이픈, 공백, 슬래시 등 무시)"""
        import re
        # 하이픈, 공백, 언더스코어, 슬래시 등을 제거하여 정규화
        message_normalized = re.sub(r'[-_\s/]+', '', message.lower())
        
        for yacht in self.yacht_data.get('yachts', []):
            yacht_name = yacht.get('name', '')
            if not yacht_name:
                continue
            
            # 요트 이름도 정규화 (슬래시도 제거)
            yacht_name_normalized = re.sub(r'[-_\s/]+', '', yacht_name.lower())
            
            # 정규화된 이름이 메시지에 포함되어 있는지 확인
            if yacht_name_normalized in message_normalized:
                return yacht_name
            
            # 부분 매칭도 시도 (예: "farr40" -> "Farr 40", "j70" -> "J/70")
            if yacht_name_normalized and message_normalized.find(yacht_name_normalized) != -1:
                return yacht_name
            
            # 숫자만 있는 경우도 시도 (예: "j70" -> "J/70")
            # 요트 이름에서 숫자 추출
            yacht_numbers = re.findall(r'\d+', yacht_name_normalized)
            message_numbers = re.findall(r'\d+', message_normalized)
            if yacht_numbers and message_numbers:
                # 숫자가 일치하고, 요트 이름의 문자 부분이 메시지에 포함되어 있으면 매칭
                if yacht_numbers[0] == message_numbers[0]:
                    yacht_letters = re.sub(r'\d+', '', yacht_name_normalized)
                    message_letters = re.sub(r'\d+', '', message_normalized)
                    if yacht_letters and yacht_letters in message_letters:
                        return yacht_name
        
        return None
    
    def _analyze_yacht_data(self, yacht_name: str, section_filter: str = None) -> str:
        """요트 데이터 종합 분석
        
        Args:
            yacht_name: 분석할 요트 이름
            section_filter: 특정 섹션만 추출 (예: "사용 목적", "적합성 평가", "정비 권장사항")
        """
        # 요트 정보 찾기
        yacht = None
        for y in self.yacht_data.get('yachts', []):
            if y.get('name', '').lower() == yacht_name.lower():
                yacht = y
                break
        
        if not yacht:
            return f"'{yacht_name}' 요트 정보를 찾을 수 없습니다."
        
        # Gemini AI를 사용한 상세 분석
        if self.has_gemini:
            # 분석 시작 메시지 출력 (즉시 표시)
            print("📊 요트 데이터를 분석 중입니다... 잠시만 기다려주세요. ⏳")
            sys.stdout.flush()  # 버퍼 강제 출력
            
            try:
                analysis_prompt = f"""다음 요트 데이터를 종합적으로 분석해주세요:

요트 정보:
{json.dumps(yacht, ensure_ascii=False, indent=2)}

부품 데이터 (해당 요트):
{json.dumps(self._get_yacht_parts(yacht_name), ensure_ascii=False, indent=2)[:2000]}

위 데이터를 바탕으로 다음을 포함한 종합 분석을 제공해주세요:
1. 요트의 주요 특징 및 스펙 요약
2. 치수 및 성능 분석
3. 부품 구성 및 정비 주기 분석
4. 사용 목적에 따른 적합성 평가
5. 관리 및 정비 권장사항

친근하고 전문적인 톤으로 답변해주세요."""
                
                response = self.model.generate_content(analysis_prompt)
                full_result = response.text
                
                # 특정 섹션만 요청된 경우 필터링
                if section_filter:
                    filtered = self._extract_section_from_analysis(full_result, section_filter)
                    if filtered:
                        result = f"📊 **{yacht_name} - {section_filter}**\n\n{filtered}"
                    else:
                        result = f"📊 **{yacht_name} 종합 분석**\n\n{full_result}\n\n💡 요청하신 '{section_filter}' 섹션을 찾지 못해 전체 분석을 보여드립니다."
                else:
                    result = f"📊 **{yacht_name} 종합 분석**\n\n{full_result}"
                
                sys.stdout.flush()  # 버퍼 강제 출력
                return result
            except Exception as e:
                # AI 분석 실패 시 기본 정보 제공
                return self._format_full_yacht_info(yacht)
        else:
            # Gemini AI 없을 때 기본 정보 제공
            return self._format_full_yacht_info(yacht)
    
    def _extract_section_from_analysis(self, full_text: str, section_keyword: str) -> str:
        """분석 결과에서 특정 섹션만 추출
        
        Args:
            full_text: 전체 분석 텍스트
            section_keyword: 찾을 섹션 키워드 (예: "사용 목적", "적합성", "정비")
        
        Returns:
            추출된 섹션 텍스트 또는 빈 문자열
        """
        import re
        
        # 섹션 제목 패턴 (###, ####, ** 등으로 시작하는 제목)
        lines = full_text.split('\n')
        section_start = -1
        section_end = len(lines)
        section_level = 0
        
        # 키워드 정규화 (공백, 특수문자 제거)
        keyword_normalized = re.sub(r'[^\w가-힣]', '', section_keyword.lower())
        
        for i, line in enumerate(lines):
            # 마크다운 제목 감지 (###, ####, **)
            if re.match(r'^#+\s+', line) or re.match(r'^\*\*.*\*\*', line):
                # 현재 줄에서 키워드 검색
                line_normalized = re.sub(r'[^\w가-힣]', '', line.lower())
                
                if keyword_normalized in line_normalized:
                    section_start = i
                    # 제목 레벨 파악 (# 개수)
                    match = re.match(r'^(#+)\s+', line)
                    section_level = len(match.group(1)) if match else 2
                    continue
                
                # 섹션 시작 후 같은 레벨의 다른 제목을 만나면 종료
                if section_start >= 0:
                    match = re.match(r'^(#+)\s+', line)
                    current_level = len(match.group(1)) if match else 2
                    if current_level <= section_level:
                        section_end = i
                        break
        
        if section_start >= 0:
            return '\n'.join(lines[section_start:section_end]).strip()
        
        return ""
    
    def _get_yacht_parts(self, yacht_name: str) -> List[Dict]:
        """요트의 부품 목록 가져오기"""
        parts_list = []
        for yacht_data in self.parts_data.get('yachts', []):
            if yacht_data.get('name', '').lower() == yacht_name.lower():
                parts_list = yacht_data.get('parts', [])
                break
        return parts_list
    
    def clear_history(self):
        """대화 히스토리 초기화"""
        self.chat_history = []
        self.current_yacht_registration = None
        print("🔄 대화 기록이 초기화되었습니다.")
    
    def get_history(self) -> List[Dict[str, str]]:
        """대화 히스토리 반환"""
        return self.chat_history
    
    def _generate_yacht_id(self, yacht_name: str) -> str:
        """
        요트 ID 생성 함수
        
        규칙:
        - 소문자 변환
        - 공백 → 하이픈 (-)
        - 슬래시 (/) → 하이픈 (-)
        - 특수문자 제거
        - 여러 개의 연속된 하이픈을 하나로 통합
        
        예시:
        - "J/70" → "j-70"
        - "OCEANIS 46.1" → "oceanis-46.1"
        - "Grand Soleil 42 Long Cruise" → "grand-soleil-42-long-cruise"
        """
        import re
        
        # 1. 소문자 변환
        yacht_id = yacht_name.lower()
        
        # 2. 슬래시를 하이픈으로 변환
        yacht_id = yacht_id.replace("/", "-")
        
        # 3. 공백을 하이픈으로 변환
        yacht_id = yacht_id.replace(" ", "-")
        
        # 4. 허용된 문자만 남기기 (영문, 숫자, 하이픈, 점)
        yacht_id = re.sub(r'[^a-z0-9\-\.]', '', yacht_id)
        
        # 5. 연속된 하이픈을 하나로 통합
        yacht_id = re.sub(r'-+', '-', yacht_id)
        
        # 6. 앞뒤 하이픈 제거
        yacht_id = yacht_id.strip('-')
        
        return yacht_id
    
    def _convert_analysis_to_registration(self, analysis_result: Dict) -> Dict:
        """분석 결과를 요트 등록 형식으로 변환 (ID 포함)"""
        doc_info = analysis_result.get("documentInfo", {})
        yacht_specs = analysis_result.get("yachtSpecs", {})
        parts = analysis_result.get("parts", [])
        
        yacht_name = doc_info.get("yachtModel") or doc_info.get("title", "Unknown Yacht")
        manufacturer = doc_info.get("manufacturer", "")
        
        # 🆕 요트 ID 생성
        yacht_id = self._generate_yacht_id(yacht_name)
        
        dimensions = yacht_specs.get("dimensions", {})
        loa_str = dimensions.get("LOA", "") or dimensions.get("loa", "")
        loa = self._parse_number(loa_str)
        beam_str = dimensions.get("Beam", "") or dimensions.get("beam", "")
        beam = self._parse_number(beam_str)
        draft_str = dimensions.get("Draft", "") or dimensions.get("draft", "")
        draft = self._parse_number(draft_str)
        displacement_str = dimensions.get("Displacement", "") or dimensions.get("displacement", "")
        displacement = self._parse_number(displacement_str)
        mast_height_str = dimensions.get("mastHeight", "") or dimensions.get("mastHeight", "")
        mast_height = self._parse_number(mast_height_str)
        
        engine = yacht_specs.get("engine", {})
        sail_area = yacht_specs.get("sailArea", {})
        
        part_list = []
        for part in parts:
            part_list.append({
                "name": part.get("name", ""),
                "manufacturer": part.get("manufacturer", ""),
                "model": part.get("model", ""),
                "interval": part.get("interval") if part.get("interval") else None,
                "latestMaintenanceDate": part.get("latestMaintenanceDate") or part.get("lastMaintenanceDate") or part.get("servicedOn") or None
            })
        
        registration_data = {
            "id": yacht_id,  # 🆕 요트 ID 추가
            "basicInfo": {
                "id": yacht_id,  # 🆕 basicInfo에도 ID 추가
                "name": yacht_name,
                "nickName": yacht_name,
                "manufacturer": manufacturer,
                "type": doc_info.get("documentType", ""),
                "year": "",
                "designer": "",
                "manual": analysis_result.get("fileInfo", {}).get("fileName", "")
            },
            "specifications": {
                "dimensions": {
                    "loa": loa,
                    "lwl": None,
                    "beam": beam,
                    "draft": draft,
                    "displacement": displacement,
                    "mastHeight": mast_height
                },
                "sailArea": {
                    "mainSailArea": self._parse_number(sail_area.get("mainsail", "")),
                    "jibSailArea": self._parse_number(sail_area.get("jib", "")),
                    "spinnakerSailArea": self._parse_number(sail_area.get("spinnaker", "")),
                    "totalSailArea": self._parse_number(sail_area.get("total", ""))
                },
                "engine": {
                    "type": engine.get("type", ""),
                    "power": engine.get("output") or engine.get("power", ""),
                    "model": engine.get("model", "")
                },
                "hull": {"hullMaterial": "", "deckMaterial": "", "keelType": ""},
                "accommodations": {"berths": None, "cabins": None, "heads": None},
                "capacity": {"fuelCapacity": None, "waterCapacity": None},
                "performance": {"maxSpeed": None, "cruisingSpeed": None},
                "ceCertification": "",
                "description": f"PDF 매뉴얼에서 자동 추출: {doc_info.get('title', '')}",
                "features": ""
            },
            "parts": part_list
        }
        
        return registration_data
    
    def _parse_number(self, value) -> Optional[float]:
        """문자열에서 숫자 추출"""
        if not value or not isinstance(value, str):
            return None
        import re
        match = re.search(r'(\d+\.?\d*)', value)
        if match:
            try:
                return float(match.group(1))
            except ValueError:
                return None
        return None
    
    def _generate_registration_completion_message(self, analysis_result: Dict, registration_data: Dict) -> str:
        """등록 완료 메시지 생성"""
        doc_info = analysis_result.get("documentInfo", {})
        yacht_name = doc_info.get("yachtModel") or doc_info.get("title", "요트")
        parts_count = len(analysis_result.get("parts", []))
        
        message = f"""✅ 등록이 완료됐습니다! 🎉

**등록된 요트 정보:**
⛵ 모델: {yacht_name}
🏭 제조사: {doc_info.get('manufacturer', 'N/A')}
📄 문서 유형: {doc_info.get('documentType', 'N/A')}

**추출된 정보:**
📏 치수 정보: {'추출됨' if registration_data['specifications']['dimensions']['loa'] else '없음'}
🔧 부품 정보: {parts_count}개 부품 추출됨
⚙️ 엔진 정보: {'추출됨' if registration_data['specifications']['engine']['type'] else '없음'}

요트가 성공적으로 등록되었습니다! 이제 부품 관리와 정비 일정을 설정할 수 있습니다.

추가로 궁금한 점이 있으시면 언제든 물어보세요! 💬"""
        
        return message
    
    def _save_registration_to_json(self, registration_data: Dict, analysis_result: Dict):
        """등록 데이터를 JSON 파일로 저장"""
        try:
            self._add_to_yacht_specifications(registration_data, analysis_result)
            self._save_to_registered_yachts(registration_data, analysis_result)
            self._save_parts_to_json_files(registration_data, analysis_result)
            print("💾 JSON 파일에 저장 완료!")
        except Exception as e:
            print(f"⚠️ JSON 파일 저장 중 오류: {e}")
    
    def _add_to_yacht_specifications(self, registration_data: Dict, analysis_result: Dict):
        """yacht_specifications.json에 요트 추가 (ID 포함)"""
        try:
            spec_file = 'data/yacht_specifications.json'
            if os.path.exists(spec_file):
                with open(spec_file, 'r', encoding='utf-8') as f:
                    data = json.load(f)
            else:
                data = {
                    "version": "1.0",
                    "description": "요트 상세 스펙 데이터베이스",
                    "lastUpdated": datetime.now().strftime("%Y-%m-%d"),
                    "yachts": []
                }
            
            basic_info = registration_data.get("basicInfo", {})
            specs = registration_data.get("specifications", {})
            
            # 🆕 registration_data에서 ID 가져오기 (없으면 생성)
            yacht_id = registration_data.get("id") or basic_info.get("id")
            if not yacht_id:
                yacht_id = self._generate_yacht_id(basic_info.get("name", ""))
            
            existing_ids = [y.get("id") for y in data.get("yachts", [])]
            if yacht_id in existing_ids:
                # 기존 요트 업데이트
                for yacht in data["yachts"]:
                    if yacht.get("id") == yacht_id:
                        yacht.update({
                            "id": yacht_id,  # 🆕 ID 명시
                            "name": basic_info.get("name", ""),
                            "manufacturer": basic_info.get("manufacturer", ""),
                            "type": basic_info.get("type", ""),
                            "manual": basic_info.get("manual", ""),
                            **self._convert_specs_to_yacht_format(specs)
                        })
                        break
            else:
                # 새 요트 추가
                new_yacht = {
                    "id": yacht_id,  # 🆕 ID 우선 배치
                    "name": basic_info.get("name", ""),
                    "manufacturer": basic_info.get("manufacturer", ""),
                    "type": basic_info.get("type", ""),
                    "manual": basic_info.get("manual", ""),
                    **self._convert_specs_to_yacht_format(specs)
                }
                data["yachts"].append(new_yacht)
            
            data["lastUpdated"] = datetime.now().strftime("%Y-%m-%d")
            with open(spec_file, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
            
            print(f"✅ {spec_file}에 저장됨 (ID: {yacht_id})")
        except Exception as e:
            print(f"⚠️ yacht_specifications.json 저장 실패: {e}")
    
    def _convert_specs_to_yacht_format(self, specs: Dict) -> Dict:
        """등록 데이터 스펙을 yacht_specifications.json 형식으로 변환"""
        dimensions = specs.get("dimensions", {})
        sail_area = specs.get("sailArea", {})
        engine = specs.get("engine", {})
        hull = specs.get("hull", {})
        accommodations = specs.get("accommodations", {})
        capacity = specs.get("capacity", {})
        performance = specs.get("performance", {})
        
        def format_dimension(value, unit="m"):
            if value is None:
                return None
            return {"value": value, "unit": unit, "display": f"{value}{unit}"}
        
        result = {}
        
        if dimensions:
            result["dimensions"] = {}
            if dimensions.get("loa"):
                result["dimensions"]["loa"] = format_dimension(dimensions.get("loa"))
            if dimensions.get("lwl"):
                result["dimensions"]["lwl"] = format_dimension(dimensions.get("lwl"))
            if dimensions.get("beam"):
                result["dimensions"]["beam"] = format_dimension(dimensions.get("beam"))
            if dimensions.get("draft"):
                result["dimensions"]["draft"] = format_dimension(dimensions.get("draft"))
            if dimensions.get("displacement"):
                result["dimensions"]["displacement"] = format_dimension(dimensions.get("displacement"), "kg")
            if dimensions.get("mastHeight"):
                result["dimensions"]["mastHeight"] = format_dimension(dimensions.get("mastHeight"))
        
        if sail_area:
            result["sailArea"] = {
                "mainSailArea": sail_area.get("mainSailArea"),
                "jibSailArea": sail_area.get("jibSailArea"),
                "spinnakerSailArea": sail_area.get("spinnakerSailArea"),
                "totalSailArea": sail_area.get("totalSailArea")
            }
        
        if engine:
            result["engine"] = {
                "type": engine.get("type", ""),
                "power": engine.get("power", ""),
                "model": engine.get("model", "")
            }
        
        if hull:
            result["hull"] = {
                "hullMaterial": hull.get("hullMaterial", ""),
                "deckMaterial": hull.get("deckMaterial", ""),
                "keelType": hull.get("keelType", "")
            }
        
        if accommodations:
            result["accommodations"] = {
                "berths": accommodations.get("berths"),
                "cabins": accommodations.get("cabins"),
                "heads": accommodations.get("heads")
            }
        
        if capacity:
            result["capacity"] = {
                "fuelCapacity": capacity.get("fuelCapacity"),
                "waterCapacity": capacity.get("waterCapacity")
            }
        
        if performance:
            result["performance"] = {
                "maxSpeed": performance.get("maxSpeed"),
                "cruisingSpeed": performance.get("cruisingSpeed")
            }
        
        if specs.get("ceCertification"):
            result["ceCertification"] = specs.get("ceCertification")
        if specs.get("description"):
            result["description"] = specs.get("description")
        if specs.get("features"):
            result["features"] = specs.get("features")
        
        return result
    
    def _save_to_registered_yachts(self, registration_data: Dict, analysis_result: Dict):
        """등록된 요트를 registered_yachts.json에 저장 (ID 포함)"""
        try:
            reg_file = 'data/registered_yachts.json'
            if os.path.exists(reg_file):
                with open(reg_file, 'r', encoding='utf-8') as f:
                    data = json.load(f)
            else:
                data = {
                    "schemaVersion": "5.0",
                    "description": "사용자가 등록한 요트 목록 (chatbot_unified.py로 등록)",
                    "lastUpdated": datetime.now().strftime("%Y-%m-%d"),
                    "totalYachts": 0,
                    "yachts": []
                }
            
            # 🆕 요트 ID 가져오기 (registration_data에 이미 포함됨)
            yacht_id = registration_data.get("id") or registration_data.get("basicInfo", {}).get("id")
            
            registration_entry = {
                "id": yacht_id,  # 🆕 최상위 ID
                "registrationDate": datetime.now().isoformat(),
                "source": "PDF Upload",
                "pdfFile": analysis_result.get("fileInfo", {}).get("fileName", ""),
                "registrationData": registration_data,
                "analysisResult": {
                    "documentInfo": analysis_result.get("documentInfo", {}),
                    "partsCount": len(analysis_result.get("parts", [])),
                    "analysisStatus": "success" if "error" not in analysis_result else "error"
                }
            }
            
            data["yachts"].append(registration_entry)
            data["totalYachts"] = len(data["yachts"])
            data["lastUpdated"] = datetime.now().strftime("%Y-%m-%d")
            
            with open(reg_file, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
            
            print(f"✅ {reg_file}에 저장됨 (ID: {yacht_id})")
        except Exception as e:
            print(f"⚠️ registered_yachts.json 저장 실패: {e}")
    
    def _save_parts_to_json_files(self, registration_data: Dict, analysis_result: Dict):
        """부품 정보를 각 JSON 파일에 저장 (요트 ID 사용)"""
        try:
            basic_info = registration_data.get("basicInfo", {})
            yacht_name = basic_info.get("name", "")
            
            # 🆕 registration_data에서 ID 가져오기
            yacht_id = registration_data.get("id") or basic_info.get("id")
            if not yacht_id:
                yacht_id = self._generate_yacht_id(yacht_name)
            
            manufacturer = basic_info.get("manufacturer", "")
            manual_pdf = basic_info.get("manual", "")
            parts = analysis_result.get("parts", [])
            
            if not parts:
                print("⚠️ 추출된 부품이 없어 부품 JSON 파일 저장을 건너뜁니다.")
                return
            
            self._add_to_yacht_parts_database(yacht_id, yacht_name, manufacturer, manual_pdf, parts)
            self._add_to_extracted_parts_detailed(yacht_id, yacht_name, manufacturer, manual_pdf, parts)
            self._add_to_extracted_parts(yacht_id, yacht_name, manufacturer, manual_pdf, parts)
            self._add_to_parts_app_data(yacht_id, yacht_name, manufacturer, manual_pdf, parts)
            
            print(f"✅ 부품 정보가 {len(parts)}개 JSON 파일에 저장됨 (Yacht ID: {yacht_id})")
        except Exception as e:
            print(f"⚠️ 부품 JSON 파일 저장 중 오류: {e}")
            import traceback
            traceback.print_exc()
    
    def _add_to_yacht_parts_database(self, yacht_id: str, yacht_name: str, manufacturer: str, manual_pdf: str, parts: List[Dict]):
        """yacht_parts_database.json에 부품 추가 (요트 ID 사용)"""
        try:
            db_file = 'data/yacht_parts_database.json'
            if os.path.exists(db_file):
                with open(db_file, 'r', encoding='utf-8') as f:
                    data = json.load(f)
            else:
                data = {"yachts": []}
            
            yacht_entry = None
            for yacht in data.get("yachts", []):
                if yacht.get("id") == yacht_id:
                    yacht_entry = yacht
                    break
            
            if not yacht_entry:
                yacht_entry = {
                    "id": yacht_id,  # 🆕 올바른 ID 사용
                    "name": yacht_name,
                    "manufacturer": manufacturer,
                    "type": "",
                    "length": None,
                    "officialWebsite": None,
                    "manualPDF": manual_pdf,
                    "dimensions": {},
                    "parts": {
                        "rigging": {"physicalParts": [], "maintenanceItems": []},
                        "sails": {"physicalParts": [], "maintenanceItems": []},
                        "engine": {"physicalParts": [], "maintenanceItems": []},
                        "hull": {"physicalParts": [], "maintenanceItems": []},
                        "electrical": {"physicalParts": [], "maintenanceItems": []},
                        "plumbing": {"physicalParts": [], "maintenanceItems": []}
                    }
                }
                data["yachts"].append(yacht_entry)
            
            parts_dict = yacht_entry.get("parts", {})
            
            for part in parts:
                category = part.get("category", "rigging").lower()
                name = part.get("name", "")
                if not name:
                    continue
                
                if category in ["rigging", "rig"]:
                    cat_key = "rigging"
                elif category in ["sails", "sail"]:
                    cat_key = "sails"
                elif category in ["engine", "motor"]:
                    cat_key = "engine"
                elif category in ["hull", "deck"]:
                    cat_key = "hull"
                elif category in ["electrical", "electric", "electronics"]:
                    cat_key = "electrical"
                elif category in ["plumbing", "water"]:
                    cat_key = "plumbing"
                else:
                    cat_key = "rigging"
                
                if cat_key not in parts_dict:
                    parts_dict[cat_key] = {"physicalParts": [], "maintenanceItems": []}
                
                physical_part = {
                    "id": f"{yacht_id}-{cat_key}-{len(parts_dict[cat_key]['physicalParts']) + 1:02d}",
                    "category": category.capitalize(),
                    "name": name,
                    "partNumber": part.get("model", ""),
                    "manufacturer": part.get("manufacturer", ""),
                    "maintenanceInterval": f"{part.get('interval', 12)}개월" if part.get("interval") else "Annual inspection"
                }
                
                parts_dict[cat_key]["physicalParts"].append(physical_part)
            
            with open(db_file, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
            
            print(f"✅ yacht_parts_database.json에 저장됨 (Yacht ID: {yacht_id})")
        except Exception as e:
            print(f"⚠️ yacht_parts_database.json 저장 실패: {e}")
    
    def _add_to_extracted_parts_detailed(self, yacht_id: str, yacht_name: str, manufacturer: str, manual_pdf: str, parts: List[Dict]):
        """extracted_yacht_parts_detailed.json에 부품 추가"""
        try:
            file_path = 'data/extracted_yacht_parts_detailed.json'
            if os.path.exists(file_path):
                with open(file_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
            else:
                data = {"yachts": []}
            
            yacht_entry = None
            for yacht in data.get("yachts", []):
                if yacht.get("id") == yacht_id:
                    yacht_entry = yacht
                    break
            
            if not yacht_entry:
                yacht_entry = {
                    "id": yacht_id,
                    "name": yacht_name,
                    "manufacturer": manufacturer,
                    "manualPDF": manual_pdf,
                    "parts": {
                        "rigging": [], "sails": [], "engine": [],
                        "hull": [], "electrical": [], "plumbing": []
                    }
                }
                data["yachts"].append(yacht_entry)
            
            parts_dict = yacht_entry.get("parts", {})
            
            for part in parts:
                category = part.get("category", "rigging").lower()
                name = part.get("name", "")
                if not name:
                    continue
                
                if category in ["rigging", "rig"]:
                    cat_key = "rigging"
                elif category in ["sails", "sail"]:
                    cat_key = "sails"
                elif category in ["engine", "motor"]:
                    cat_key = "engine"
                elif category in ["hull", "deck"]:
                    cat_key = "hull"
                elif category in ["electrical", "electric", "electronics"]:
                    cat_key = "electrical"
                elif category in ["plumbing", "water"]:
                    cat_key = "plumbing"
                else:
                    cat_key = "rigging"
                
                if cat_key not in parts_dict:
                    parts_dict[cat_key] = []
                
                part_entry = {
                    "name": name,
                    "description": f"{manufacturer} {yacht_name} - {name}",
                    "specifications": [
                        part.get("model", ""),
                        part.get("manufacturer", ""),
                        f"Interval: {part.get('interval', 'N/A')} months" if part.get("interval") else ""
                    ]
                }
                
                parts_dict[cat_key].append(part_entry)
            
            with open(file_path, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
            
            print(f"✅ extracted_yacht_parts_detailed.json에 저장됨")
        except Exception as e:
            print(f"⚠️ extracted_yacht_parts_detailed.json 저장 실패: {e}")
    
    def _add_to_extracted_parts(self, yacht_id: str, yacht_name: str, manufacturer: str, manual_pdf: str, parts: List[Dict]):
        """extracted_yacht_parts.json에 부품 추가 (간단한 형식)"""
        try:
            file_path = 'data/extracted_yacht_parts.json'
            if os.path.exists(file_path):
                with open(file_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
            else:
                data = {"yachts": []}
            
            # data가 리스트인 경우 딕셔너리로 변환
            if isinstance(data, list):
                data = {"yachts": data}
            
            # data가 딕셔너리가 아니거나 "yachts" 키가 없는 경우
            if not isinstance(data, dict):
                data = {"yachts": []}
            elif "yachts" not in data:
                data["yachts"] = []
            
            yacht_entry = None
            for yacht in data.get("yachts", []):
                if yacht.get("id") == yacht_id:
                    yacht_entry = yacht
                    break
            
            if not yacht_entry:
                yacht_entry = {
                    "id": yacht_id,
                    "name": yacht_name,
                    "manufacturer": manufacturer,
                    "parts": []
                }
                if not isinstance(data.get("yachts"), list):
                    data["yachts"] = []
                data["yachts"].append(yacht_entry)
            
            # 기존 부품 목록 가져오기 (중복 방지)
            if not isinstance(yacht_entry, dict):
                yacht_entry = {"parts": []}
            existing_parts = yacht_entry.get("parts", [])
            existing_part_names = {p.get("name", "") for p in existing_parts if isinstance(p, dict)}
            
            for part in parts:
                # parts가 dict인지 확인
                if not isinstance(part, dict):
                    continue
                    
                name = part.get("name", "")
                if not name or name in existing_part_names:
                    continue
                
                part_entry = {
                    "name": name,
                    "manufacturer": part.get("manufacturer", ""),
                    "model": part.get("model", ""),
                    "category": part.get("category", "rigging"),
                    "interval": part.get("interval")
                }
                
                yacht_entry["parts"].append(part_entry)
                existing_part_names.add(name)
            
            with open(file_path, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
            
            print(f"✅ extracted_yacht_parts.json에 저장됨")
        except Exception as e:
            print(f"⚠️ extracted_yacht_parts.json 저장 실패: {e}")
            import traceback
            traceback.print_exc()
    
    def _add_to_parts_app_data(self, yacht_id: str, yacht_name: str, manufacturer: str, manual_pdf: str, parts: List[Dict]):
        """yacht_parts_app_data.json에 부품 추가"""
        try:
            file_path = 'data/yacht_parts_app_data.json'
            if os.path.exists(file_path):
                with open(file_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
            else:
                data = {"yachts": []}
            
            yacht_entry = None
            for yacht in data.get("yachts", []):
                if yacht.get("id") == yacht_id:
                    yacht_entry = yacht
                    break
            
            if not yacht_entry:
                yacht_entry = {
                    "id": yacht_id,
                    "name": yacht_name,
                    "manufacturer": manufacturer,
                    "parts": []
                }
                data["yachts"].append(yacht_entry)
            
            for part in parts:
                name = part.get("name", "")
                if not name:
                    continue
                
                part_entry = {
                    "name": name,
                    "manufacturer": part.get("manufacturer", ""),
                    "model": part.get("model", ""),
                    "category": part.get("category", "rigging"),
                    "maintenanceInterval": part.get("interval", 12)
                }
                
                yacht_entry["parts"].append(part_entry)
            
            with open(file_path, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
            
            print(f"✅ yacht_parts_app_data.json에 저장됨")
        except Exception as e:
            print(f"⚠️ yacht_parts_app_data.json 저장 실패: {e}")
    
    def get_registration_data(self) -> Optional[Dict]:
        """현재 등록 중인 요트 데이터 반환"""
        return self.current_yacht_registration
    
    def _analyze_document_for_preview(self, file_path: str, extracted_text: str) -> Dict:
        """문서 분석 - 부품 목록만 빠르게 추출 (미리보기용)"""
        if len(extracted_text) > 30000:
            extracted_text = extracted_text[:30000] + "\n\n[... 텍스트 일부만 분석 ...]"
        
        if not self.has_gemini:
            return {"error": "문서 분석 기능은 Gemini API가 필요합니다."}
        
        # 부품 목록만 빠르게 추출하는 프롬프트
        prompt = f"""다음은 요트 또는 부품 매뉴얼에서 추출한 텍스트입니다:

{extracted_text}

---

## 📋 작업 지시사항

**목적: 부품 목록 빠른 추출 (사용자 선택을 위한 미리보기)**

매뉴얼에서 발견되는 **모든 부품**을 간단히 나열하세요.

**부품 추출 규칙:**

1. **ID 생성**: `part-{{순번}}` (예: part-001, part-002)
2. **필수 필드**: id, name, category
3. **선택 필드**: manufacturer, model, interval (있으면 포함)
4. **중복 통합**: 같은 이름의 부품은 하나로 통합
5. **세부 부품 포함**: 
   - 엔진의 구성 부품 (예: 해수펌프, 연료여과기)
   - 시스템의 하위 부품 (예: 배터리, 충전기)
   - 모두 개별 부품으로 추출

**중요도 표시:**
- importance: "high" (핵심 부품, 엔진/돛/키 등)
- importance: "medium" (중요 하위 부품, 펌프/필터 등)
- importance: "low" (소모품, 볼트/개스킷 등)

**응답 형식:**

```json
{{
  "documentInfo": {{
    "title": "문서 제목",
    "yachtModel": "요트 모델명 (있다면)",
    "manufacturer": "제조사",
    "documentType": "Owner's Manual / Parts Manual / Service Manual"
  }},
  "parts": [
    {{
      "id": "part-001",
      "name": "부품명",
      "category": "Engine / Rigging / Deck / Electrical / Plumbing",
      "manufacturer": "제조사 (있다면)",
      "model": "모델명 (있다면)",
      "interval": 12,
      "importance": "high",
      "description": "간단한 설명 (선택사항)"
    }}
  ],
  "totalPartsCount": 59,
  "summary": {{
    "high": 5,
    "medium": 20,
    "low": 34
  }}
}}
```

**주의사항:**
- 모든 부품을 빠짐없이 추출 (개수 제한 없음)
- 부품이 많아도 (50개 이상) 모두 나열
- 불확실한 정보는 null로 표시
- JSON 형식으로만 응답"""
        
        try:
            print("🤖 부품 목록 추출 중...")
            response = self.model.generate_content(prompt)
            result_text = response.text
            
            # JSON 추출
            if "```json" in result_text:
                json_start = result_text.find("```json") + 7
                json_end = result_text.find("```", json_start)
                result_text = result_text[json_start:json_end].strip()
            elif "```" in result_text:
                json_start = result_text.find("```") + 3
                json_end = result_text.find("```", json_start)
                result_text = result_text[json_start:json_end].strip()
            
            result = json.loads(result_text)
            result["fileInfo"] = {
                "fileName": os.path.basename(file_path),
                "filePath": file_path
            }
            
            print(f"✅ {len(result.get('parts', []))}개 부품 발견!")
            return result
            
        except Exception as e:
            print(f"❌ 부품 목록 추출 실패: {e}")
            return {"error": str(e)}
    
    def _analyze_with_selected_parts(self, file_path: str, extracted_text: str, selected_part_ids: List[str]) -> Dict:
        """선택된 부품만 포함하여 전체 문서 분석"""
        if len(extracted_text) > 30000:
            extracted_text = extracted_text[:30000] + "\n\n[... 텍스트 일부만 분석 ...]"
        
        if not self.has_gemini:
            return {"error": "문서 분석 기능은 Gemini API가 필요합니다."}
        
        # 선택된 부품 ID 목록을 문자열로 변환
        selected_ids_str = ", ".join(selected_part_ids) if selected_part_ids else "없음"
        
        prompt = f"""다음은 요트 매뉴얼에서 추출한 텍스트입니다:

{extracted_text}

---

## 📋 작업 지시사항 (선택된 부품만 추출)

**선택된 부품 ID:** {selected_ids_str}

위 ID에 해당하는 부품만 상세하게 분석하세요. 다른 부품은 무시하세요.

**부품이 선택되지 않은 경우 (ID: 없음):**
- parts 배열을 빈 배열 []로 반환
- 요트 기본 정보(documentInfo, yachtSpecs)는 여전히 추출

**응답 형식 (Schema 5.0):**

```json
{{
  "schemaVersion": "5.0",
  "analyzedAt": "2025-11-28T...",
  "documentInfo": {{
    "title": "...",
    "yachtModel": "...",
    "manufacturer": "...",
    "documentType": "..."
  }},
  "yachtSpecs": {{
    "standard": {{
      "dimensions": {{}},
      "engine": {{}},
      "sailArea": {{}}
    }},
    "additional": {{}}
  }},
  "parts": [
    // 선택된 ID의 부품만 포함 (상세 정보)
    {{
      "id": "선택된 ID",
      "name": "...",
      "manufacturer": "...",
      "model": "...",
      "category": "...",
      "interval": 12,
      "specifications": {{}},
      "maintenanceDetails": {{}}
    }}
  ]
}}
```

**중요:**
- 선택되지 않은 부품은 절대 포함하지 마세요
- JSON 형식으로만 응답"""
        
        try:
            print(f"🤖 선택된 {len(selected_part_ids)}개 부품 상세 분석 중...")
            response = self.model.generate_content(prompt)
            result_text = response.text
            
            # JSON 추출
            if "```json" in result_text:
                json_start = result_text.find("```json") + 7
                json_end = result_text.find("```", json_start)
                result_text = result_text[json_start:json_end].strip()
            elif "```" in result_text:
                json_start = result_text.find("```") + 3
                json_end = result_text.find("```", json_start)
                result_text = result_text[json_start:json_end].strip()
            
            result = json.loads(result_text)
            result["fileInfo"] = {
                "fileName": os.path.basename(file_path),
                "filePath": file_path
            }
            
            print(f"✅ {len(result.get('parts', []))}개 부품 분석 완료!")
            return result
            
        except Exception as e:
            print(f"❌ 부품 분석 실패: {e}")
            return {"error": str(e)}


def run_interactive_mode(api_key: str = None):
    """대화형 모드 실행"""
    chatbot = UnifiedYachtChatbot(api_key=api_key, mode="interactive")
    
    print("\n💡 사용 팁:")
    print("  - 자연스럽게 질문하세요 (예: 'Farr 40 크기 알려줘')")
    print("  - PDF 파일 경로를 입력하면 자동으로 분석합니다")
    print("  - '/list' - 요트 목록 보기")
    print("  - '/info' - 데이터 정보 보기")
    print("  - '/help' - 도움말 보기")
    print("  - '/quit' 또는 '/exit' - 종료")
    print("\n" + "=" * 60 + "\n")
    
    while True:
        try:
            user_input = input("👤 You: ").strip()
            
            if not user_input:
                continue
            
            if user_input.lower() in ['/quit', '/exit', '/q']:
                print("\n👋 HooAah Yacht 챗봇을 종료합니다. 안녕히 가세요!")
                break
            
            if user_input.lower() == '/clear':
                chatbot.clear_history()
                continue
            
            print("\n🤖 AI: ", end="", flush=True)
            response = chatbot.chat(user_input)
            print(response)
            print()
            
        except KeyboardInterrupt:
            print("\n\n👋 프로그램을 종료합니다.")
            break
        except Exception as e:
            print(f"\n❌ 오류 발생: {e}")


def run_api_server(api_key: str = None, port: int = 5000):
    """API 서버 모드 실행"""
    if not HAS_FLASK:
        print("❌ Flask가 설치되지 않았습니다. pip install flask flask-cors")
        return
    
    chatbot = UnifiedYachtChatbot(api_key=api_key, mode="api")
    
    app = Flask(__name__)
    CORS(app)
    
    chatbot_sessions = {}
    
    def get_or_create_chatbot(session_id: str):
        if session_id not in chatbot_sessions:
            chatbot_sessions[session_id] = UnifiedYachtChatbot(api_key=api_key, mode="api")
        return chatbot_sessions[session_id]
    
    @app.route('/api/chat', methods=['POST'])
    def chat():
        try:
            data = request.get_json()
            if not data or 'message' not in data:
                return jsonify({"success": False, "error": "메시지가 필요합니다."}), 400
            
            user_message = data['message']
            session_id = data.get('session_id', 'default')
            pdf_file_path = data.get('pdf_file_path')  # 모바일 앱에서 파일 경로 전달
            
            chatbot = get_or_create_chatbot(session_id)
            ai_response = chatbot.chat(user_message, pdf_file_path=pdf_file_path)
            
            return jsonify({
                "success": True,
                "response": ai_response,
                "session_id": session_id,
                "timestamp": datetime.now().isoformat()
            })
        except Exception as e:
            return jsonify({"success": False, "error": str(e)}), 500
    
    @app.route('/api/chat/upload', methods=['POST'])
    def upload_pdf():
        """
        PDF 파일 업로드 API (모바일 앱용)
        
        Request:
        - multipart/form-data
        - file: PDF 파일
        - message: 사용자 메시지 (선택사항)
        - session_id: 세션 ID (선택사항)
        """
        try:
            if 'file' not in request.files:
                return jsonify({"success": False, "error": "파일이 필요합니다."}), 400
            
            file = request.files['file']
            if file.filename == '':
                return jsonify({"success": False, "error": "파일이 선택되지 않았습니다."}), 400
            
            if not file.filename.lower().endswith('.pdf'):
                return jsonify({"success": False, "error": "PDF 파일만 업로드 가능합니다."}), 400
            
            # 파일 저장
            upload_folder = 'uploads'
            os.makedirs(upload_folder, exist_ok=True)
            
            if secure_filename:
                filename = secure_filename(file.filename)
            else:
                # secure_filename이 없으면 기본 파일명 사용
                filename = file.filename.replace(' ', '_')
            file_path = os.path.join(upload_folder, filename)
            file.save(file_path)
            
            # 세션 정보
            session_id = request.form.get('session_id', 'default')
            user_message = request.form.get('message', f'PDF 파일 업로드: {filename}')
            
            # 챗봇으로 처리
            chatbot = get_or_create_chatbot(session_id)
            ai_response = chatbot.chat(user_message, pdf_file_path=file_path)
            
            return jsonify({
                "success": True,
                "response": ai_response,
                "session_id": session_id,
                "file_name": filename,
                "timestamp": datetime.now().isoformat()
            })
        except Exception as e:
            return jsonify({"success": False, "error": str(e)}), 500
    
    @app.route('/api/yacht/register', methods=['POST'])
    def register_yacht():
        """
        요트 PDF 등록 API - JSON 형식으로 추출 데이터 반환
        
        Request:
        - multipart/form-data
        - file: PDF 파일
        
        Response:
        - JSON 형식의 추출된 요트 데이터 (자연어 없음)
        """
        try:
            if 'file' not in request.files:
                return jsonify({"success": False, "error": "파일이 필요합니다."}), 400
            
            file = request.files['file']
            if file.filename == '':
                return jsonify({"success": False, "error": "파일이 선택되지 않았습니다."}), 400
            
            if not file.filename.lower().endswith('.pdf'):
                return jsonify({"success": False, "error": "PDF 파일만 업로드 가능합니다."}), 400
            
            # 파일 저장
            upload_folder = 'uploads'
            os.makedirs(upload_folder, exist_ok=True)
            
            if secure_filename:
                filename = secure_filename(file.filename)
            else:
                filename = file.filename.replace(' ', '_')
            file_path = os.path.join(upload_folder, filename)
            file.save(file_path)
            
            # 세션 정보
            session_id = request.form.get('session_id', 'default')
            
            # 챗봇 인스턴스 가져오기
            chatbot = get_or_create_chatbot(session_id)
            
            # 텍스트 추출
            print(f"📄 파일 분석 시작: {filename}", flush=True)
            extracted_text = chatbot._extract_text_from_file(file_path)
            
            if not extracted_text or len(extracted_text.strip()) < 100:
                return jsonify({
                    "success": False,
                    "error": f"{filename}에서 텍스트를 추출할 수 없습니다."
                }), 400
            
            # AI 분석 (JSON 형식으로)
            print(f"🤖 AI 분석 중...", flush=True)
            analysis_result = chatbot._analyze_document_directly(file_path, extracted_text)
            
            # 분석 실패 확인
            if "error" in analysis_result:
                return jsonify({
                    "success": False,
                    "error": analysis_result.get("error", "분석 실패")
                }), 500
            
            # 등록 데이터 변환
            registration_data = chatbot._convert_analysis_to_registration(analysis_result)
            
            # JSON 파일 저장
            chatbot._save_registration_to_json(registration_data, analysis_result)
            
            # JSON 형식으로 응답 (자연어 없음)
            return jsonify({
                "success": True,
                "fileName": filename,
                "timestamp": datetime.now().isoformat(),
                "yacht": {
                    "basicInfo": registration_data.get("basicInfo", {}),
                    "specifications": registration_data.get("specifications", {}),
                    "parts": registration_data.get("parts", [])
                },
                "analysisResult": {
                    "documentInfo": analysis_result.get("documentInfo", {}),
                    "yachtSpecs": analysis_result.get("yachtSpecs", {}),
                    "detailedDimensions": analysis_result.get("detailedDimensions", {}),
                    "exterior": analysis_result.get("exterior", {}),
                    "groundTackle": analysis_result.get("groundTackle", {}),
                    "sailInventory": analysis_result.get("sailInventory", []),
                    "deckEquipment": analysis_result.get("deckEquipment", {}),
                    "accommodations": analysis_result.get("accommodations", {}),
                    "tanks": analysis_result.get("tanks", {}),
                    "electricalSystem": analysis_result.get("electricalSystem", {}),
                    "electronics": analysis_result.get("electronics", {}),
                    "plumbingSystem": analysis_result.get("plumbingSystem", {}),
                    "parts": analysis_result.get("parts", []),
                    "maintenance": analysis_result.get("maintenance", [])
                }
            }), 200
            
        except Exception as e:
            print(f"❌ Error: {e}")
            import traceback
            traceback.print_exc()
            return jsonify({"success": False, "error": str(e)}), 500
    
    @app.route('/api/yacht/preview-parts', methods=['POST'])
    def preview_parts():
        """
        PDF 부품 미리보기 API - 추출 가능한 부품 목록만 반환
        
        Request:
        - multipart/form-data
        - file: PDF 파일
        
        Response:
        - 부품 목록 (사용자가 선택할 수 있도록)
        """
        try:
            if 'file' not in request.files:
                return jsonify({"success": False, "error": "파일이 필요합니다."}), 400
            
            file = request.files['file']
            if file.filename == '':
                return jsonify({"success": False, "error": "파일이 선택되지 않았습니다."}), 400
            
            if not file.filename.lower().endswith('.pdf'):
                return jsonify({"success": False, "error": "PDF 파일만 업로드 가능합니다."}), 400
            
            # 파일 저장
            upload_folder = 'uploads'
            os.makedirs(upload_folder, exist_ok=True)
            
            if secure_filename:
                filename = secure_filename(file.filename)
            else:
                filename = file.filename.replace(' ', '_')
            file_path = os.path.join(upload_folder, filename)
            file.save(file_path)
            
            # 세션 정보
            session_id = request.form.get('session_id', 'default')
            chatbot = get_or_create_chatbot(session_id)
            
            # 텍스트 추출
            print(f"📄 파일 분석 시작: {filename}", flush=True)
            extracted_text = chatbot._extract_text_from_file(file_path)
            
            if not extracted_text or len(extracted_text.strip()) < 100:
                return jsonify({
                    "success": False,
                    "error": f"{filename}에서 텍스트를 추출할 수 없습니다."
                }), 400
            
            # 부품 목록만 빠르게 추출
            preview_result = chatbot._analyze_document_for_preview(file_path, extracted_text)
            
            if "error" in preview_result:
                return jsonify({
                    "success": False,
                    "error": preview_result.get("error", "분석 실패")
                }), 500
            
            # 응답
            return jsonify({
                "success": True,
                "fileName": filename,
                "filePath": file_path,  # 다음 단계에서 사용
                "timestamp": datetime.now().isoformat(),
                "documentInfo": preview_result.get("documentInfo", {}),
                "parts": preview_result.get("parts", []),
                "totalPartsCount": preview_result.get("totalPartsCount", len(preview_result.get("parts", []))),
                "summary": preview_result.get("summary", {}),
                "message": f"총 {len(preview_result.get('parts', []))}개의 부품이 발견되었습니다. 등록할 부품을 선택하세요."
            }), 200
            
        except Exception as e:
            print(f"❌ Error: {e}")
            import traceback
            traceback.print_exc()
            return jsonify({"success": False, "error": str(e)}), 500
    
    @app.route('/api/yacht/register-selected', methods=['POST'])
    def register_selected_parts():
        """
        선택된 부품만 등록 API
        
        Request (JSON):
        {
          "filePath": "uploads/manual.pdf",
          "selectedPartIds": ["part-001", "part-005", "part-010"],
          "session_id": "default"
        }
        
        Response:
        - 선택된 부품만 포함된 전체 분석 결과
        """
        try:
            data = request.get_json()
            
            if not data:
                return jsonify({"success": False, "error": "JSON 데이터가 필요합니다."}), 400
            
            file_path = data.get('filePath')
            selected_part_ids = data.get('selectedPartIds', [])
            session_id = data.get('session_id', 'default')
            
            if not file_path:
                return jsonify({"success": False, "error": "filePath가 필요합니다."}), 400
            
            if not os.path.exists(file_path):
                return jsonify({"success": False, "error": "파일을 찾을 수 없습니다."}), 404
            
            # 챗봇 인스턴스
            chatbot = get_or_create_chatbot(session_id)
            
            # 텍스트 재추출
            print(f"📄 파일 재분석: {os.path.basename(file_path)}", flush=True)
            extracted_text = chatbot._extract_text_from_file(file_path)
            
            if not extracted_text or len(extracted_text.strip()) < 100:
                return jsonify({
                    "success": False,
                    "error": "텍스트 추출 실패"
                }), 400
            
            # 선택된 부품만 상세 분석
            analysis_result = chatbot._analyze_with_selected_parts(file_path, extracted_text, selected_part_ids)
            
            if "error" in analysis_result:
                return jsonify({
                    "success": False,
                    "error": analysis_result.get("error", "분석 실패")
                }), 500
            
            # 등록 데이터 변환
            registration_data = chatbot._convert_analysis_to_registration(analysis_result)
            
            # JSON 파일 저장
            chatbot._save_registration_to_json(registration_data, analysis_result)
            
            # 응답
            return jsonify({
                "success": True,
                "fileName": os.path.basename(file_path),
                "timestamp": datetime.now().isoformat(),
                "selectedPartsCount": len(selected_part_ids),
                "registeredPartsCount": len(analysis_result.get("parts", [])),
                "yacht": {
                    "basicInfo": registration_data.get("basicInfo", {}),
                    "specifications": registration_data.get("specifications", {}),
                    "parts": registration_data.get("parts", [])
                },
                "analysisResult": {
                    "documentInfo": analysis_result.get("documentInfo", {}),
                    "yachtSpecs": analysis_result.get("yachtSpecs", {}),
                    "parts": analysis_result.get("parts", [])
                },
                "message": f"{len(analysis_result.get('parts', []))}개 부품이 등록되었습니다."
            }), 200
            
        except Exception as e:
            print(f"❌ Error: {e}")
            import traceback
            traceback.print_exc()
            return jsonify({"success": False, "error": str(e)}), 500
    
    @app.route('/api/chat/history', methods=['GET'])
    def get_history():
        try:
            session_id = request.args.get('session_id', 'default')
            chatbot = get_or_create_chatbot(session_id)
            history = chatbot.get_history()
            return jsonify({"success": True, "history": history})
        except Exception as e:
            return jsonify({"success": False, "error": str(e)}), 500
    
    @app.route('/api/yacht/analyze', methods=['GET'])
    def analyze_yacht_by_name():
        """요트 이름으로 부품 정보 조회 (Backend 연동용)"""
        try:
            yacht_name = request.args.get('yacht_name', '').strip()
            
            if not yacht_name:
                return jsonify({
                    "success": False,
                    "error": "yacht_name parameter is required"
                }), 400
            
            # 챗봇 인스턴스 생성
            chatbot = get_or_create_chatbot('backend-api')
            
            # 요트 데이터 조회
            yacht_data = None
            for yacht in chatbot.yacht_data.get('yachts', []):
                if yacht.get('name', '').lower() == yacht_name.lower():
                    yacht_data = yacht
                    break
            
            if not yacht_data:
                return jsonify({
                    "success": False,
                    "error": f"Yacht '{yacht_name}' not found"
                }), 404
            
            # 부품 데이터 조회
            yacht_id = yacht_data.get('id', '')
            parts_list = []
            
            for yacht_parts in chatbot.parts_data.get('yachts', []):
                if yacht_parts.get('id') == yacht_id:
                    parts_list = yacht_parts.get('parts', [])
                    break
            
            # Backend DTO 형식으로 변환
            parts_dto = []
            for part in parts_list:
                parts_dto.append({
                    "id": part.get('id', ''),
                    "name": part.get('name', ''),
                    "manufacturer": part.get('manufacturer', ''),
                    "model": part.get('model', ''),
                    "interval": part.get('interval'),
                    "maintenanceDetails": {
                        "recommendedInterval": part.get('maintenanceDetails', {}).get('recommendedInterval', ''),
                        "maintenanceMethod": part.get('maintenanceDetails', {}).get('maintenanceMethod', ''),
                        "notes": part.get('maintenanceDetails', {}).get('notes', '')
                    }
                })
            
            return jsonify({
                "success": True,
                "yachtId": yacht_id,
                "yachtName": yacht_data.get('name', ''),
                "parts": parts_dto,
                "totalParts": len(parts_dto)
            }), 200
            
        except Exception as e:
            print(f"❌ Error in analyze_yacht_by_name: {e}")
            import traceback
            traceback.print_exc()
            return jsonify({
                "success": False,
                "error": str(e)
            }), 500
    
    @app.route('/api/yacht/analyze-pdf', methods=['POST'])
    def analyze_pdf_file():
        """PDF 파일로 부품 정보 분석 (Backend 연동용)"""
        try:
            if 'file' not in request.files:
                return jsonify({
                    "success": False,
                    "error": "No file provided"
                }), 400
            
            file = request.files['file']
            
            if file.filename == '':
                return jsonify({
                    "success": False,
                    "error": "Empty filename"
                }), 400
            
            if not file.filename.lower().endswith('.pdf'):
                return jsonify({
                    "success": False,
                    "error": "Only PDF files are supported"
                }), 400
            
            # 임시 파일로 저장
            import tempfile
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
                file.save(tmp_file.name)
                tmp_path = tmp_file.name
            
            try:
                # 챗봇 인스턴스로 분석
                chatbot = get_or_create_chatbot('backend-api')
                
                # PDF 텍스트 추출
                extracted_text = chatbot._extract_text_from_pdf(tmp_path)
                
                if not extracted_text or len(extracted_text.strip()) < 100:
                    return jsonify({
                        "success": False,
                        "error": "Unable to extract text from PDF"
                    }), 400
                
                # AI 분석
                analysis_result = chatbot._analyze_document_directly(tmp_path, extracted_text)
                
                if not analysis_result or "error" in analysis_result:
                    return jsonify({
                        "success": False,
                        "error": analysis_result.get("error", "Analysis failed")
                    }), 500
                
                # Backend DTO 형식으로 변환
                parts_list = analysis_result.get('parts', [])
                parts_dto = []
                
                for part in parts_list:
                    parts_dto.append({
                        "id": part.get('id', ''),
                        "name": part.get('name', ''),
                        "manufacturer": part.get('manufacturer', ''),
                        "model": part.get('model', ''),
                        "interval": part.get('maintenanceDetails', {}).get('interval'),
                        "maintenanceDetails": {
                            "recommendedInterval": part.get('maintenanceDetails', {}).get('recommendedInterval', ''),
                            "maintenanceMethod": part.get('maintenanceDetails', {}).get('maintenanceMethod', ''),
                            "notes": part.get('maintenanceDetails', {}).get('notes', '')
                        }
                    })
                
                yacht_name = analysis_result.get('documentInfo', {}).get('yachtName', 'Unknown')
                yacht_id = chatbot._generate_yacht_id(yacht_name)
                
                return jsonify({
                    "success": True,
                    "yachtId": yacht_id,
                    "yachtName": yacht_name,
                    "parts": parts_dto,
                    "totalParts": len(parts_dto),
                    "documentInfo": {
                        "fileName": file.filename,
                        "manufacturer": analysis_result.get('documentInfo', {}).get('manufacturer', ''),
                        "model": analysis_result.get('documentInfo', {}).get('model', ''),
                        "year": analysis_result.get('documentInfo', {}).get('year')
                    }
                }), 200
                
            finally:
                # 임시 파일 삭제
                import os
                if os.path.exists(tmp_path):
                    os.remove(tmp_path)
                    
        except Exception as e:
            print(f"❌ Error in analyze_pdf_file: {e}")
            import traceback
            traceback.print_exc()
            return jsonify({
                "success": False,
                "error": str(e)
            }), 500
    
    @app.route('/api/health', methods=['GET'])
    def health_check():
        """헬스체크 엔드포인트"""
        try:
            chatbot = get_or_create_chatbot('health-check')
            yacht_count = len(chatbot.yacht_data.get('yachts', []))
            
            return jsonify({
                "status": "healthy",
                "timestamp": datetime.now().isoformat(),
                "yachtCount": yacht_count,
                "version": "5.0"
            }), 200
        except Exception as e:
            return jsonify({
                "status": "unhealthy",
                "error": str(e),
                "timestamp": datetime.now().isoformat()
            }), 500
    
    print("=" * 60)
    print("🌐 HooAah Yacht AI Chatbot API Server")
    print("=" * 60)
    print(f"🚀 서버 시작: http://localhost:{port}")
    print("📡 API 엔드포인트:")
    print("  [챗봇용]")
    print("  - POST /api/chat - 채팅 메시지 전송")
    print("  - POST /api/chat/upload - PDF 업로드 (자연어 응답)")
    print("  - POST /api/yacht/register - 요트 등록 (JSON 응답)")
    print("  - GET /api/chat/history - 대화 기록 조회")
    print()
    print("  [부품 선택 등록] ⭐ NEW")
    print("  - POST /api/yacht/preview-parts - PDF 부품 미리보기 (선택 전)")
    print("  - POST /api/yacht/register-selected - 선택된 부품만 등록")
    print()
    print("  [Backend 연동용]")
    print("  - GET /api/yacht/analyze?yacht_name={name} - 요트 이름으로 부품 조회")
    print("  - POST /api/yacht/analyze-pdf - PDF 파일 분석")
    print("  - GET /api/health - 서버 상태 확인")
    print("=" * 60)
    print()
    
    app.run(host='0.0.0.0', port=port, debug=True)


def main():
    """메인 함수"""
    # 시작 메시지 (즉시 표시)
    print("🚀 HooAah Yacht 챗봇을 시작하는 중...", flush=True)
    print("⏳ 잠시만 기다려주세요...", flush=True)
    print()
    
    parser = argparse.ArgumentParser(description='HooAah Yacht 통합 챗봇')
    parser.add_argument('--mode', choices=['interactive', 'api'], default='interactive',
                        help='실행 모드 (interactive: 대화형, api: API 서버)')
    parser.add_argument('--api-key', type=str, help='Gemini API 키 (선택사항, .env 파일 또는 환경변수 사용)')
    parser.add_argument('--port', type=int, default=5000, help='API 서버 포트 (기본: 5000)')
    
    args = parser.parse_args()
    
    # API 키 우선순위: 명령줄 인자 > 환경변수 (.env 파일)
    api_key = args.api_key or os.getenv('GEMINI_API_KEY')
    if not api_key:
        print("❌ 오류: GEMINI_API_KEY가 설정되지 않았습니다.")
        print("💡 다음 중 하나를 선택하세요:")
        print("   1. .env 파일 생성: .env.example을 참고하여 .env 파일을 만들고 GEMINI_API_KEY를 설정")
        print("   2. 환경변수 설정: export GEMINI_API_KEY=your-api-key (Linux/Mac) 또는 set GEMINI_API_KEY=your-api-key (Windows)")
        print("   3. 명령줄 옵션: --api-key YOUR_API_KEY")
        sys.exit(1)
    
    if args.mode == 'api':
        run_api_server(api_key=api_key, port=args.port)
    else:
        run_interactive_mode(api_key=api_key)


if __name__ == "__main__":
    main()

